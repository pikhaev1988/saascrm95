from __future__ import annotations

from dataclasses import dataclass
from io import BytesIO
from statistics import mean
import re
from collections import Counter, defaultdict
from datetime import date

from django.db.models import Avg, Count, Min, Max, Q

from exams.models import EgePassingThreshold, ExamResult, TaskResult
from exams.passing import gve_subject_label, is_gve_exam
from organizations.models import School
from users.ai import enhance_exam_analysis
from users.task_topics import (
    default_grade_label,
    default_skill_profile,
    domain_from_subject_key,
    is_expanded_answer_task,
    part2_start_task,
    subject_key as resolve_subject_key,
    topic_for_task,
)


@dataclass
class ExamData:
    subject: str
    date: str
    students_count: int
    avg_score: float
    min_score: float
    max_score: float
    pass_rate: float
    tasks: list[dict]
    strong_tasks: list[int]
    weak_tasks: list[int]
    recommendations: list[str]
    topic_deficits: list[str]
    exam_type: str
    score_values: list[float]
    exam_year: int
    dynamics: list[dict]


STOPWORDS = {
    "кл",
    "класс",
    "пункт",
    "раздел",
    "тема",
    "задание",
    "элемент",
    "содержание",
    "по",
    "и",
    "в",
    "на",
    "с",
    "для",
    "из",
    "как",
    "при",
    "к",
}

# Лёгкие subject-профили задают только стартовые акценты.
# Основная модель компетенций ниже строится динамически от доменов и паттернов результатов.
SUBJECT_SKILL_PROFILES = {
    "math": ["многошаговое рассуждение", "моделирование", "стратегия вычислений"],
    "russian": ["текстовая интерпретация", "языковой анализ", "письменная аргументация"],
    "biology": ["системный анализ", "процессная интерпретация", "экспериментальное мышление"],
    "chemistry": ["расчетная стратегия", "анализ реакций", "причинно-следственные связи"],
    "physics": ["построение модели", "анализ закономерностей", "интерпретация результата"],
    "social": ["анализ кейса", "понятийный аппарат", "аргументированное объяснение"],
    "history": ["хронология", "источниковедческий анализ", "историческая аргументация"],
    "informatics": ["алгоритмическое мышление", "логический анализ", "программная стратегия"],
    "literature": ["интерпретация текста", "сопоставительный анализ", "аргументация"],
    "geography": ["пространственный анализ", "анализ данных", "практико-ориентированная интерпретация"],
    "foreign": ["понимание текста", "языковая точность", "коммуникативная аргументация"],
}

CONTROL_ACTIONS_BY_DOMAIN = {
    "math": [
        "практикум по разбору многошаговых задач и декомпозиции условий",
        "тренинг доказательного решения с разбором типичных ошибок",
        "серия задач на моделирование прикладных ситуаций",
    ],
    "russian": [
        "орфографические и пунктуационные тренажеры по дефицитным правилам",
        "практика синтаксического разбора с комментированием решения",
        "регулярные мини-сочинения по критериям ЕГЭ с экспертной обратной связью",
    ],
    "biology": [
        "работа со схемами, рисунками и таблицами по проблемной теме",
        "генетические и биологические задачи с пошаговым разбором",
        "моделирование процессов и экосистем с обсуждением причин ошибок",
    ],
    "chemistry": [
        "разбор расчетных задач с отработкой единиц измерения и алгоритмов",
        "практика составления уравнений реакций в разных классах веществ",
        "тематические лабораторные мини-кейсы по условиям протекания реакций",
    ],
    "physics": [
        "практикум по выбору физической модели и закона для задачи",
        "серия задач на качественный анализ явлений и графическую интерпретацию",
        "тренинг расчетных задач с проверкой размерности и логики ответа",
    ],
    "history": [
        "хронологические тренинги по периодам и ключевым событиям",
        "практика анализа исторических источников и сопоставления позиций",
        "упражнения на причинно-следственные связи и историческую аргументацию",
    ],
    "social": [
        "разбор кейсов по экономике, праву и политике с понятийным аппаратом",
        "практика аргументированных ответов по общественным ситуациям",
        "тренинг заданий второй части по критериям оценивания",
    ],
    "informatics": [
        "алгоритмические тренировки с декомпозицией и псевдокодом",
        "практика программирования и отладки типовых задач",
        "разбор логических и табличных задач с формализацией данных",
    ],
    "literature": [
        "практика аналитического чтения и выделения авторской позиции",
        "сопоставительный анализ фрагментов и художественных средств",
        "тренинг письменных ответов по критериям литературного анализа",
    ],
    "geography": [
        "картографические практикумы по чтению и интерпретации карт",
        "разбор задач на географические закономерности и факторы размещения",
        "тренинг заданий с использованием статистических и пространственных данных",
    ],
    "foreign": [
        "лексико-грамматические практикумы по дефицитным темам",
        "тренинг чтения/аудирования с выделением ключевой информации",
        "отработка письменных заданий в формате ЕГЭ по критериям",
    ],
}


def _subject_key(subject_name: str, exam_type: str = "ege") -> str:
    return resolve_subject_key(subject_name, exam_type)


def _domain_from_subject_key(subject_key: str) -> str:
    return domain_from_subject_key(subject_key)


def _detect_grading_system(subject_key: str, exam_type: str = "ege", max_score: float | None = None) -> str:
    if (exam_type or "").lower() == "oge":
        return "5-point"
    if max_score is not None and max_score <= 5:
        return "5-point"
    if subject_key == "math_basic":
        return "5-point"
    return "100-point"


def _grading_model_label_ru(grading_system: str) -> str:
    return "пятибалльная (оценки 2–5)" if grading_system == "5-point" else "стобалльная (тестовые баллы 0–100)"


def _exam_type_label_ru(exam_type: str) -> str:
    return "ОГЭ" if str(exam_type or "").strip().lower() == "oge" else "ЕГЭ"


def _risk_level_ru_from_rates(success_or_pass_rate: float, low: float = 40, mid: float = 60) -> str:
    """Уровень риска для отображения пользователю (русские подписи)."""
    if success_or_pass_rate < low:
        return "Высокий"
    if success_or_pass_rate < mid:
        return "Средний"
    return "Низкий"


def _strip_curriculum_codes(topic: str) -> str:
    text = (topic or "").strip()
    text = re.sub(r"\b\d{1,2}\s*кл\.,?\s*п\.\s*[\d.]+\s*", "", text, flags=re.I)
    text = re.sub(r"\bп\.\s*[\d.]+\s*", "", text, flags=re.I)
    text = re.sub(r"\b\d{2,}\.\d{1,2}\.\d{1,2}\.?\d*\b", "", text)
    text = re.sub(r"\b\d+\.\d+\.\d+\b", "", text)
    text = re.sub(r"\s{2,}", " ", text)
    return text.strip(" ;,.")


def _normalize_topic(topic: str) -> str:
    clean = _strip_curriculum_codes(topic)
    clean = re.sub(r"\b(10|11)\s*класс\b", "", clean, flags=re.I)
    clean = re.sub(
        r"\b(контроль|усилить|развивать|формирование|отработка|совершенствование|разбор|работа по|повторение)\b",
        "",
        clean,
        flags=re.I,
    )
    clean = re.sub(r"\b(умение|навык|заданий?|темы?)\b", "", clean, flags=re.I)
    clean = re.sub(r"\s{2,}", " ", clean).strip(" ;,.:-")
    if not clean:
        return "Тематический блок по спецификации"
    chunks = [x.strip(" ;,.") for x in re.split(r"[;•,.](?=\s*[А-ЯA-Z])", clean) if x.strip()]
    candidate = chunks[0] if chunks else clean
    lower = candidate.lower()
    canonical = {
        r"биология как наука|методы изучения живой природы|методы исследован": "Методы биологических исследований",
        r"человек\s*[-–]\s*биосоциальный|анатом|физиолог|биология человека": "Биология человека",
        r"живые системы|общая биолог|уровни организации жив": "Общая биология",
        r"планиметр": "Планиметрия",
        r"стереометр|пространствен": "Стереометрия",
        r"уравнен|неравен": "Уравнения и неравенства",
        r"алгебр|вычисл|преобразован": "Алгебраические вычисления",
        r"график|таблиц|диаграм": "Графики и таблицы",
        r"пунктуац.*сложн": "Пунктуация сложного предложения",
        r"орфограф": "Орфография",
        r"синтакс": "Синтаксис",
        r"эколог|экосистем|биосфер": "Экология",
        r"генет|наследств": "Генетика",
    }
    for pattern, label in canonical.items():
        if re.search(pattern, lower):
            return label
    candidate = re.sub(r"\b(5|6|7|8|9|10|11)\s*класс(ы)?\b", "", candidate, flags=re.I).strip(" ;,.-")
    candidate = candidate[:1].upper() + candidate[1:] if candidate else candidate
    return candidate


def _fallback_domain_for_oge(topic: str, task_id: int | None = None) -> str:
    text = (topic or "").lower()
    if any(token in text for token in ("пунктуац", "знак препин")):
        return "Пунктуация"
    if any(token in text for token in ("орфограф", "правопис")):
        return "Орфография"
    if any(token in text for token in ("синтакс", "предложен")):
        return "Синтаксис"
    if any(token in text for token in ("текст", "чтени", "автор")):
        return "Работа с текстом"
    if any(token in text for token in ("алгебр", "уравнен", "вычисл")):
        return "Алгебра"
    if any(token in text for token in ("геометр", "планиметр", "стереометр")):
        return "Геометрия"
    if any(token in text for token in ("вероятн", "статист", "комбинатор")):
        return "Вероятность и статистика"
    if any(token in text for token in ("алгоритм", "программ", "информат")):
        return "Алгоритмы"
    if task_id is not None:
        return f"Тематический блок задания №{task_id}"
    return "Тематический блок"


def _topic_tokens(text: str) -> list[str]:
    words = re.findall(r"[А-Яа-яA-Za-z]+", (text or "").lower())
    return [w for w in words if len(w) > 3 and w not in STOPWORDS]


def _semantic_key(topic: str) -> str:
    tokens = _topic_tokens(topic)
    if not tokens:
        return _normalize_topic(topic).lower()
    base = tokens[:3]
    return " ".join(base)


def _merge_semantic_domains(topic_blocks: dict[str, list[str]]) -> dict[str, list[str]]:
    merged: dict[str, list[str]] = {}
    signatures: dict[str, str] = {}
    for block, topics in topic_blocks.items():
        sign = _semantic_key(block)
        if sign in signatures:
            target = signatures[sign]
            for topic in topics:
                if topic not in merged[target]:
                    merged[target].append(topic)
        else:
            signatures[sign] = block
            merged[block] = list(topics)
    return merged


def _score_groups(scores: list[float], subject_key: str, exam_type: str = "ege") -> dict[str, int]:
    # Для OГЭ используем 5-балльную модель при шкале 2-5.
    if (exam_type == "oge" and scores and max(scores) <= 5) or subject_key == "math_basic":
        groups = {"grade_2": 0, "grade_3": 0, "grade_4": 0, "grade_5": 0}
        for score in scores:
            grade = int(round(score))
            if grade <= 2:
                groups["grade_2"] += 1
            elif grade == 3:
                groups["grade_3"] += 1
            elif grade == 4:
                groups["grade_4"] += 1
            else:
                groups["grade_5"] += 1
        return groups

    groups = {"weak": 0, "basic": 0, "good": 0, "high": 0}
    for score in scores:
        if score <= 35:
            groups["weak"] += 1
        elif score <= 60:
            groups["basic"] += 1
        elif score <= 80:
            groups["good"] += 1
        else:
            groups["high"] += 1
    return groups


def _group_topic_deficits(topics: list[str], subject_key: str) -> dict[str, list[str]]:
    grouped: dict[str, list[str]] = defaultdict(list)
    labels = {}
    for topic in topics:
        readable = _normalize_topic(topic)
        key = _semantic_key(readable)
        if key not in labels:
            labels[key] = readable
        if readable not in grouped[key]:
            grouped[key].append(readable)
    sorted_items = sorted(grouped.items(), key=lambda kv: len(kv[1]), reverse=True)
    return {labels[key]: values for key, values in sorted_items if values}


def _extract_classes(topic: str) -> list[str]:
    text = topic or ""
    classes = set()
    for m in re.finditer(r"(\d+)\s*кл", text, flags=re.I):
        g = int(m.group(1))
        if 1 <= g <= 11:
            classes.add(str(g))
    for m in re.finditer(r"(\d+)\s*[–-]\s*(\d+)\s*кл", text, flags=re.I):
        a, b = sorted((int(m.group(1)), int(m.group(2))))
        for g in range(a, b + 1):
            if 1 <= g <= 11:
                classes.add(str(g))
    return sorted(classes, key=lambda x: int(x))


def _dominant_group(groups: dict[str, int]) -> str:
    if not groups or sum(groups.values()) == 0:
        return "не определена"
    key = max(groups.items(), key=lambda item: item[1])[0]
    labels = {
        "weak": "слабый уровень",
        "basic": "базовый уровень",
        "good": "хороший уровень",
        "high": "высокий уровень",
        "grade_2": "2 (неудовлетворительно)",
        "grade_3": "3 (удовлетворительно)",
        "grade_4": "4 (хорошо)",
        "grade_5": "5 (отлично)",
    }
    return labels.get(key, "не определена")


def _skill_lines_for_subject(subject_key: str) -> list[str]:
    # Базовый fallback; детальный вывод строится функцией _infer_competencies().
    domain = _domain_from_subject_key(subject_key)
    profile = SUBJECT_SKILL_PROFILES.get(domain) or []
    if profile:
        return profile
    return ["аналитическое мышление", "аргументация решения", "применение знаний в новой ситуации"]


def _infer_competencies(subject_key: str, domains: list[str], weak_candidates: list[dict], second_part_avg: float, first_part_avg: float) -> list[str]:
    inferred: list[str] = []
    domain_text = " ".join(domains).lower()
    domain = _domain_from_subject_key(subject_key)
    rules_by_domain = {
        "russian": [
            (r"пунктуац", "синтаксический анализ предложения и пунктуационное моделирование"),
            (r"орфограф", "применение орфографических правил в контекстном языковом анализе"),
            (r"текст|смыслов|автор", "анализ авторской позиции и семантическая интерпретация текста"),
            (r"сочинен|аргумент", "доказательная аргументация в письменном ответе"),
            (r"синтакс", "структурный анализ синтаксических конструкций"),
        ],
        "biology": [
            (r"генет|наслед", "интерпретация схем наследования и анализ генетических закономерностей"),
            (r"эколог|экосистем|биосфер", "анализ экологических связей и устойчивости экосистем"),
            (r"анатом|физиолог|биология человека", "интерпретация физиологических процессов и функциональных взаимосвязей"),
            (r"клет|организм|систем", "анализ биологических систем и межуровневых связей"),
        ],
        "math": [
            (r"график|таблиц|диаграм", "интерпретация графических моделей и количественных зависимостей"),
            (r"геометр|стереометр|планиметр", "пространственное геометрическое моделирование и доказательное рассуждение"),
            (r"уравнен|алгебр|вычисл|функц", "алгоритмизация решения и выбор вычислительной стратегии"),
            (r"вероятн|комбинатор", "вероятностное рассуждение и комбинаторный анализ"),
        ],
        "history": [
            (r"источник|документ", "интерпретация исторического источника и вывод в контексте эпохи"),
            (r"хронолог|истори", "хронологический анализ и причинно-следственная реконструкция событий"),
        ],
        "informatics": [
            (r"алгоритм|программ", "алгоритмическая декомпозиция и формализация решения"),
            (r"логик|данн", "логическая верификация и анализ структур данных"),
        ],
    }
    rules = rules_by_domain.get(domain, [])
    for pattern, label in rules:
        if re.search(pattern, domain_text):
            inferred.append(label)
    subject_specific = {
        "biology": [
            "интерпретация биологических процессов на уровне систем и органов",
            "анализ экспериментальных данных и биологических моделей",
        ],
        "math": [
            "структурная декомпозиция условия и выбор метода решения",
            "обоснование математических переходов в многошаговом решении",
        ],
        "russian": [
            "анализ авторской позиции и смысловых акцентов текста",
            "синтаксическое моделирование и структурный анализ предложения",
            "семантическая интерпретация микротем и логики текста",
            "доказательная аргументация в развернутом письменном ответе",
            "контекстуальный орфографический анализ языкового материала",
        ],
        "history": [
            "сопоставление фактов и позиций на основе источников",
            "историческая аргументация с учетом хронологического контекста",
        ],
        "informatics": [
            "построение алгоритма и проверка граничных случаев",
            "логическая верификация решения по формальным условиям",
        ],
    }
    inferred.extend(subject_specific.get(domain, []))
    if second_part_avg and first_part_avg and second_part_avg + 10 < first_part_avg:
        inferred.append("аргументация решения в заданиях повышенного уровня сложности")
    if any(item.get("success_rate", 100) < 30 for item in weak_candidates):
        inferred.append("устойчивость базовых предметных операций в экзаменационном формате")
    if not inferred:
        inferred = _skill_lines_for_subject(subject_key)
    # de-dup preserve order
    out = []
    for item in inferred:
        if item not in out:
            out.append(item)
    return out[:4]


def _control_actions_for_subject(subject_key: str, block: str) -> list[str]:
    domain = _domain_from_subject_key(subject_key)
    actions = CONTROL_ACTIONS_BY_DOMAIN.get(domain)
    block_l = (block or "").lower()
    # Точечная адаптация действий по типу дефицита (не только по предмету).
    if domain == "math":
        if any(token in block_l for token in ("геометр", "стереометр", "планиметр")):
            return [
                "практика построений и разбор геометрических конфигураций",
                "тренинг пространственного воображения с 3D-моделями",
                "отработка задач на доказательство и обоснование каждого шага",
            ]
        if any(token in block_l for token in ("уравнен", "неравен", "функц")):
            return [
                "алгоритмы пошагового решения уравнений/неравенств",
                "обязательная проверка корней и анализ типовых ошибок",
                "серия задач на декомпозицию и выбор метода решения",
            ]
    if domain == "biology":
        if any(token in block_l for token in ("генет", "наслед")):
            return [
                "генетические задачи с пошаговым комментированием",
                "тренинг схем наследования и вероятностных расчетов",
                "разбор ошибок в записи генотипов и фенотипов",
            ]
        if any(token in block_l for token in ("эколог", "экосистем", "биосфер")):
            return [
                "моделирование экосистем и пищевых связей",
                "разбор экологических кейсов по факторам среды",
                "тренинг интерпретации биосферных процессов по диаграммам",
            ]
    if domain == "russian":
        if any(token in block_l for token in ("орфограф", "пунктуац")):
            return [
                "ежедневные орфографические и пунктуационные тренажеры",
                "разбор правил на контекстных примерах из ЕГЭ",
                "короткие диагностические диктанты с адресной коррекцией",
            ]
        if any(token in block_l for token in ("сочинен", "текст", "аргумент")):
            return [
                "практика мини-сочинений по критериям ЕГЭ",
                "разбор типовых текстов с выделением авторской позиции",
                "тренинг аргументации и логики письменного ответа",
            ]
    if domain == "history" and any(token in block_l for token in ("источник", "документ")):
        return [
            "сравнительный анализ исторических источников",
            "упражнения на атрибуцию документа и контекст эпохи",
            "тренинг исторической аргументации на основе источника",
        ]
    if actions:
        return actions[:]
    return [
        f"тематический практикум по блоку «{block}»",
        "разбор типичных ошибок с последующей коррекционной отработкой",
        "краткие диагностические срезы с анализом прогресса",
    ]


def _generate_actions(block: str, competencies: list[str], severity: str, subject_key: str) -> list[str]:
    base = _control_actions_for_subject(subject_key, block)
    block_l = (block or "").lower()
    subject_domain = _domain_from_subject_key(subject_key)
    variability_pool = [
        "мастерская по разбору типовых ловушек в формулировках заданий",
        "мини-практикум с пошаговым комментированием решений",
        "разбор эталонных и ошибочных ответов по критериям оценивания",
        "парная взаимопроверка с чек-листом предметных критериев",
        "сессия моделирования задания с устным проговариванием стратегии",
        "тренинг визуализации условия через схемы, таблицы и опорные модели",
        "структурная декомпозиция сложных задач на управляемые шаги",
        "экзаменационная симуляция с последующим аналитическим дебрифингом",
        "критериальный практикум по качеству развернутого ответа",
    ]
    if any(token in block_l for token in ("геометр", "стереометр", "пространств")):
        variability_pool.append("тренинг пространственной визуализации с построением вспомогательных чертежей")
    if any(token in block_l for token in ("генет", "наслед")):
        variability_pool.append("моделирование генетических скрещиваний с анализом вероятностных исходов")
    if subject_domain == "russian":
        variability_pool.append("практика синтаксической реконструкции предложения с пунктуационным комментарием")
    if subject_domain == "history":
        variability_pool.append("семинар по интерпретации исторического источника в контексте эпохи")
    shift = sum(ord(ch) for ch in (block or "")) % len(variability_pool)
    rotating = [variability_pool[(shift + i) % len(variability_pool)] for i in range(2)]
    adaptive = []
    if competencies:
        adaptive.append(f"целевая отработка компетенции: {competencies[0]}")
    if severity == "критический":
        adaptive.extend(
            [
                "коррекционный интенсив с повторной диагностикой после каждого цикла",
                "индивидуальные маршруты исправления типовых ошибок с фиксацией прогресса",
            ]
        )
    elif severity == "значимый":
        adaptive.append("серия тематических практикумов с промежуточной проверкой переносимости навыка")
    else:
        adaptive.append("поддерживающая тренировка на закрепление навыка в новых форматах заданий")
    combined = base + rotating + adaptive
    # unique
    result = []
    for item in combined:
        if item not in result:
            result.append(item)
    return result[:5]


def _readiness_level(
    avg_score: float,
    dominant_group: str,
    critical_count: int,
    second_part_avg: float | None,
    systemic_zero_tasks: list[int],
    grading_system: str,
    pass_rate: float = 0.0,
) -> str:
    has_part2 = second_part_avg is not None
    part2_value = second_part_avg if has_part2 else 100.0
    if grading_system == "5-point":
        if pass_rate >= 90 and avg_score >= 3.8 and critical_count == 0:
            return "высокая"
        if critical_count >= 2 or len(systemic_zero_tasks) >= 2 or (has_part2 and part2_value < 30):
            return "недостаточная"
        if "2 (" in dominant_group or avg_score < 3.2:
            return "нестабильная"
        if avg_score < 3.5 or pass_rate < 75:
            return "удовлетворительная"
        if avg_score >= 4.0 and pass_rate >= 85:
            return "устойчивая"
        return "частично устойчивая"
    if pass_rate >= 90 and avg_score >= 75 and critical_count == 0:
        return "высокая"
    if critical_count >= 3 or len(systemic_zero_tasks) >= 2 or (has_part2 and part2_value < 25):
        return "недостаточная"
    if "слабый" in dominant_group or avg_score < 55:
        return "нестабильная"
    if avg_score < 65 or pass_rate < 75:
        return "удовлетворительная"
    if avg_score >= 75 and pass_rate >= 85:
        return "устойчивая"
    return "частично устойчивая"


def _task_severity(success_rate: float) -> str:
    if success_rate < 30:
        return "критический"
    if success_rate < 50:
        return "значимый"
    if success_rate < 70:
        return "умеренный"
    return "рабочий"


def collect_subject_data_for_export(
    school_id: int,
    exam_type: str,
    subject: str,
    year: int | None,
) -> ExamData | None:
    from analytics.engine import AnalyticsEngine

    if not year:
        return None
    engine = AnalyticsEngine()
    result = engine.analyze_subject(school_id, exam_type, subject, int(year))
    if not result.valid:
        return None

    tasks = [
        {
            "id": task.task_number,
            "success_rate": task.success_rate,
            "correct": task.correct,
            "wrong": task.wrong,
            "total": task.total,
        }
        for task in result.tasks
    ]
    from analytics.engine.attempts import filter_latest_exam_results

    score_values = list(
        filter_latest_exam_results(
            ExamResult.objects.filter(
                student__school_id=school_id,
                exam__exam_type=result.exam_type,
                exam__subject=result.subject,
                exam__year=year,
            )
        ).values_list("score", flat=True)
    )
    score_values = [float(v or 0) for v in score_values]

    data = ExamData(
        subject=result.subject,
        date=result.exam_date,
        students_count=result.students_count,
        avg_score=result.avg_score,
        min_score=result.min_score,
        max_score=result.max_score,
        pass_rate=result.pass_rate,
        tasks=tasks,
        strong_tasks=result.strong_tasks,
        weak_tasks=result.weak_tasks,
        recommendations=result.recommendations,
        topic_deficits=[t.topic for t in result.topics if t.success_rate < result.avg_score][:8],
        exam_type=result.exam_type,
        score_values=score_values,
        exam_year=result.exam_year or int(year),
        dynamics=result.dynamics,
    )
    data.engine_result = result  # type: ignore[attr-defined]
    return data


def collect_exam_data_for_export(school_id: int, exam_id: int) -> ExamData | None:
    from analytics.engine import AnalyticsEngine

    engine = AnalyticsEngine()
    result = engine.analyze_exam(school_id, exam_id)
    if not result.valid:
        return None

    tasks = [
        {
            "id": task.task_number,
            "success_rate": task.success_rate,
            "correct": task.correct,
            "wrong": task.wrong,
            "total": task.total,
        }
        for task in result.tasks
    ]
    score_values = list(
        ExamResult.objects.filter(student__school_id=school_id, exam_id=exam_id).values_list("score", flat=True)
    )
    score_values = [float(v or 0) for v in score_values]

    data = ExamData(
        subject=result.subject,
        date=result.exam_date,
        students_count=result.students_count,
        avg_score=result.avg_score,
        min_score=result.min_score,
        max_score=result.max_score,
        pass_rate=result.pass_rate,
        tasks=tasks,
        strong_tasks=result.strong_tasks,
        weak_tasks=result.weak_tasks,
        recommendations=result.recommendations,
        topic_deficits=[t.topic for t in result.topics if t.success_rate < result.avg_score][:8],
        exam_type=result.exam_type,
        score_values=score_values,
        exam_year=result.exam_year or 0,
        dynamics=result.dynamics,
    )
    data.engine_result = result  # type: ignore[attr-defined]
    return data


def _collect_exam_data_legacy(school_id: int, exam_id: int) -> ExamData | None:
    results = ExamResult.objects.filter(student__school_id=school_id, exam_id=exam_id).select_related("exam")
    if not results.exists():
        return None
    first = results.first()
    exam = first.exam
    exam_type = (exam.exam_type or "ege").lower()
    exam_year = int(exam.year or exam.exam_date.year)
    stats = results.aggregate(avg=Avg("score"), min=Min("score"), max=Max("score"), cnt=Count("id"))
    students_count = int(stats["cnt"] or 0)
    pass_rate = round((results.filter(passed=True).count() / students_count) * 100, 1) if students_count else 0.0
    score_values = [float(value or 0) for value in results.values_list("score", flat=True)]

    def _is_success_token(token_value: str) -> bool:
        token = str(token_value or "").strip()
        if not token:
            return False
        if token == "+":
            return True
        if token in {"-", "0"}:
            return False
        if token.isdigit():
            return int(token) > 0
        return False

    raw_task_values = list(
        TaskResult.objects.filter(student__school_id=school_id, exam_id=exam_id)
        .values("task_number", "value")
        .order_by("task_number")
    )
    task_agg: dict[int, dict[str, int]] = {}
    for row in raw_task_values:
        task_num = int(row["task_number"])
        bucket = task_agg.setdefault(task_num, {"total": 0, "plus": 0, "minus": 0})
        bucket["total"] += 1
        if _is_success_token(row["value"]):
            bucket["plus"] += 1
        else:
            bucket["minus"] += 1

    tasks = []
    for task_num in sorted(task_agg):
        row = task_agg[task_num]
        total = int(row["total"] or 0)
        plus = int(row["plus"] or 0)
        minus = int(row["minus"] or 0)
        success_rate = round((plus / total) * 100, 1) if total else 0.0
        tasks.append(
            {
                "id": task_num,
                "success_rate": success_rate,
                "correct": plus,
                "wrong": minus,
                "total": total,
            }
        )

    strong_tasks = [item["id"] for item in tasks if item["success_rate"] >= 100.0]
    if not strong_tasks:
        strong_tasks = [item["id"] for item in sorted(tasks, key=lambda x: (-x["success_rate"], x["wrong"]))[:3]]
    weak_tasks = [item["id"] for item in tasks if item["success_rate"] < 50.0]

    topic_deficits = []
    for task_number in weak_tasks:
        topic = topic_for_task(exam.subject, task_number, exam_type)
        if topic and topic not in topic_deficits:
            topic_deficits.append(topic)

    recommendations = []
    if pass_rate < 70:
        recommendations.append("Усилить адресную работу с группой риска, проводить еженедельные мини-срезы.")
    if weak_tasks:
        weak_label = ", ".join(f"№{task}" for task in weak_tasks[:5])
        recommendations.append(f"Приоритезировать отработку заданий {weak_label} через тематические практикумы.")
    if students_count >= 10:
        if score_values and (max(score_values) - min(score_values) >= 30):
            recommendations.append("Организовать уровневые группы подготовки для выравнивания результатов.")
    if not recommendations:
        recommendations.append("Сохранить текущую стратегию подготовки и поддерживать регулярный мониторинг.")

    dynamics_years = []
    if exam_year >= 2025:
        dynamics_years = [y for y in (2023, 2024, 2025, 2026) if y <= exam_year][-4:]
    elif exam_year == 2024:
        dynamics_years = [2023, 2024]
    dynamics = []
    if dynamics_years:
        dynamics_qs = (
            ExamResult.objects.filter(
                student__school_id=school_id,
                exam__exam_type=exam_type,
                exam__subject=exam.subject,
                exam__year__in=dynamics_years,
            )
            .values("exam__year")
            .annotate(
                students=Count("id"),
                avg_score=Avg("score"),
                passed=Count("id", filter=Q(passed=True)),
            )
            .order_by("exam__year")
        )
        for row in dynamics_qs:
            students = int(row["students"] or 0)
            pass_rate_y = round((int(row["passed"] or 0) / students) * 100, 1) if students else 0.0
            dynamics.append(
                {
                    "year": int(row["exam__year"]),
                    "students": students,
                    "avg_score": round(float(row["avg_score"] or 0), 2),
                    "pass_rate": pass_rate_y,
                }
            )

    return ExamData(
        subject=exam.subject,
        date=exam.exam_date.strftime("%d.%m.%Y"),
        students_count=students_count,
        avg_score=round(float(stats["avg"] or 0), 2),
        min_score=round(float(stats["min"] or 0), 2),
        max_score=round(float(stats["max"] or 0), 2),
        pass_rate=pass_rate,
        tasks=tasks,
        strong_tasks=strong_tasks,
        weak_tasks=weak_tasks,
        recommendations=recommendations,
        topic_deficits=topic_deficits,
        exam_type=exam_type,
        score_values=score_values,
        exam_year=exam_year,
        dynamics=dynamics,
    )


def _build_analysis_payload(data: ExamData) -> dict:
    engine_result = getattr(data, "engine_result", None)
    if engine_result is not None and getattr(engine_result, "valid", False):
        from analytics.engine.adapters import to_legacy_payload

        payload = to_legacy_payload(engine_result)
        try:
            ai_result = enhance_exam_analysis(
                {
                    "exam_type": data.exam_type,
                    "subject": data.subject,
                    "students_count": data.students_count,
                    "avg_score": data.avg_score,
                    "pass_rate": data.pass_rate,
                    "draft_recommendations": payload.get("recommendations"),
                    "draft_executive_summary": payload.get("sections", {}).get("Краткие выводы", []),
                }
            )
            if isinstance(ai_result, dict):
                exec_lines = ai_result.get("executive_summary")
                if isinstance(exec_lines, list) and exec_lines:
                    payload["sections"]["Краткие выводы"] = [str(x) for x in exec_lines[:6]]
        except Exception:
            pass
        return payload

    subject_key = _subject_key(data.subject, data.exam_type)
    grading_system = _detect_grading_system(subject_key, data.exam_type, data.max_score)
    if data.exam_type == "oge":
        grading_system = "5-point"
    task_success_values = [task["success_rate"] for task in data.tasks]
    avg_task_success = round(mean(task_success_values), 1) if task_success_values else 0.0
    groups = _score_groups(data.score_values, subject_key, data.exam_type)
    part2_boundary = part2_start_task(data.exam_type, subject_key)
    first_part = [item["success_rate"] for item in data.tasks if item["id"] < part2_boundary]
    second_part = [item["success_rate"] for item in data.tasks if item["id"] >= part2_boundary]
    first_part_avg = round(mean(first_part), 1) if first_part else 0.0
    second_part_avg = round(mean(second_part), 1) if second_part else None
    systemic_zero_tasks = [item["id"] for item in data.tasks if item["success_rate"] == 0.0]
    low_tasks = [item["id"] for item in data.tasks if item["success_rate"] < 50.0]
    dominant_group = _dominant_group(groups)

    if grading_system == "5-point":
        if data.avg_score < 3:
            level = "неудовлетворительный"
        elif data.avg_score <= 3.5:
            level = "удовлетворительный"
        elif data.avg_score < 5:
            level = "хороший"
        else:
            level = "отличный"
    else:
        level = "высокий"
        if data.avg_score < 50:
            level = "низкий"
        elif data.avg_score < 65:
            level = "базовый"
        elif data.avg_score < 80:
            level = "достаточный"

    readiness = "высокая"
    if data.pass_rate < 50 or level in {"низкий", "неудовлетворительный"}:
        readiness = "недостаточная"
    elif data.pass_rate < 75:
        readiness = "частичная"

    task_topic_map: dict[int, str] = {}
    for task in data.tasks:
        task_topic_map[task["id"]] = topic_for_task(data.subject, task["id"], data.exam_type)

    strong_candidates = [t for t in data.tasks if t["success_rate"] >= 80]
    if not strong_candidates:
        strong_candidates = sorted(data.tasks, key=lambda x: (-x["success_rate"], x["wrong"]))[:3]
    weak_candidates = [t for t in data.tasks if t["success_rate"] < 50]
    relative_weak = sorted(data.tasks, key=lambda x: (x["success_rate"], -x["wrong"]))[:3]
    medium_candidates = [t for t in data.tasks if 50 <= t["success_rate"] < 80]

    strong_task_lines = []
    for task in strong_candidates[:6]:
        topic = _normalize_topic(task_topic_map.get(task["id"], ""))
        block = _normalize_topic(topic)
        skill = _skill_lines_for_subject(subject_key)[0]
        strong_task_lines.append(
            f"Задание №{task['id']} ({task['success_rate']}%): освоен блок «{block}», сформирован навык — {skill}."
        )

    weak_topics = [task_topic_map.get(task["id"], "") for task in weak_candidates if task_topic_map.get(task["id"], "")]
    topic_blocks = _merge_semantic_domains(_group_topic_deficits(weak_topics, subject_key))

    group_lines = []
    if grading_system == "5-point":
        labels = {
            "grade_2": "2 (неудовлетворительно)",
            "grade_3": "3 (удовлетворительно)",
            "grade_4": "4 (хорошо)",
            "grade_5": "5 (отлично)",
        }
        order = ("grade_2", "grade_3", "grade_4", "grade_5")
    else:
        labels = {
            "weak": "Слабый (0-35)",
            "basic": "Базовый (36-60)",
            "good": "Хороший (61-80)",
            "high": "Высокий (81-100)",
        }
        order = ("weak", "basic", "good", "high")

    for key in order:
        count = groups[key]
        if count > 0:
            share = round((count / max(data.students_count, 1)) * 100, 1)
            group_lines.append(f"{labels[key]}: {count} чел. ({share}%).")
    if not group_lines:
        group_lines.append("Недостаточно данных для группировки обучающихся.")

    severity_counter = Counter(_task_severity(task["success_rate"]) for task in data.tasks)
    critical_tasks = [f"№{task['id']}" for task in data.tasks if _task_severity(task["success_rate"]) == "критический"]
    significant_tasks = [f"№{task['id']}" for task in data.tasks if _task_severity(task["success_rate"]) == "значимый"]
    moderate_tasks = [f"№{task['id']}" for task in data.tasks if _task_severity(task["success_rate"]) == "умеренный"]

    task_lines = [
        f"Средняя успешность по заданиям: {avg_task_success}%.",
    ]
    if second_part_avg is not None:
        part2_comment = (
            "Снижение во 2 части указывает на дефицит развёрнутого решения и аргументации."
            if second_part_avg + 10 < first_part_avg
            else "Разрыв между частями не критичен."
        )
        task_lines.append(
            f"Часть 1 (задания 1–{part2_boundary - 1}): {first_part_avg}%, "
            f"часть 2 (задания {part2_boundary}+): {second_part_avg}%. {part2_comment}"
        )
    else:
        task_lines.append(
            f"Часть 1 (задания 1–{part2_boundary - 1}): {first_part_avg}%. "
            f"Данные по части 2 (задания {part2_boundary}+) в протоколе отсутствуют."
        )
    task_lines.extend(
        [
            f"Сильные задания (успешность не ниже 80%): {', '.join(f'№{x}' for x in [t['id'] for t in strong_candidates]) or 'нет'}.",
            f"Задания среднего уровня (50–80%): {', '.join(f'№{x}' for x in [t['id'] for t in medium_candidates]) or 'нет'}.",
            f"Слабые задания (успешность ниже 50%): {', '.join(f'№{x}' for x in [t['id'] for t in weak_candidates]) or 'нет'}.",
        ]
    )
    if not weak_candidates and relative_weak:
        task_lines.append(
            "Относительно более сложные задания (при отсутствии критических провалов): "
            + ", ".join(f"№{t['id']} ({t['success_rate']}%)" for t in relative_weak)
            + "."
        )
    if systemic_zero_tasks:
        task_lines.append(
            f"Системный провал: задания с 0% выполнения — {', '.join(f'№{x}' for x in systemic_zero_tasks)}."
        )
    task_lines.extend(
        [
            f"Критический уровень дефицита (<30%): {severity_counter['критический']} заданий ({', '.join(critical_tasks) if critical_tasks else 'нет'}).",
            f"Значимый уровень дефицита (30–50%): {severity_counter['значимый']} заданий ({', '.join(significant_tasks) if significant_tasks else 'нет'}).",
            f"Умеренный уровень дефицита (50–70%): {severity_counter['умеренный']} заданий ({', '.join(moderate_tasks) if moderate_tasks else 'нет'}).",
        ]
    )
    if data.exam_type == "ege":
        task_lines.append(
            "Фокус ЕГЭ на части 2: потери баллов определяются качеством аналитического рассуждения, "
            "аргументации и соответствием критериям оценивания развернутого ответа."
        )

    topic_lines = []
    if topic_blocks:
        for block, topics in topic_blocks.items():
            clean_topics = [t for t in topics if str(t).strip()]
            safe_block = str(block).strip()
            if not safe_block and data.exam_type == "oge":
                source_topic = task_topic_map.get(int(weak_candidates[0]["id"]), "") if weak_candidates else ""
                safe_block = _fallback_domain_for_oge(source_topic)
            elif not safe_block:
                safe_block = "Тематический блок"
            affected = [
                str(task["id"])
                for task in weak_candidates
                if _normalize_topic(task_topic_map.get(task["id"], "")) in topics
            ]
            if clean_topics:
                topic_lines.append(
                    f"{safe_block}: дефицитные темы — {'; '.join(clean_topics[:3])}. Задания: {', '.join('№' + x for x in affected) if affected else 'нет'}."
                )
            else:
                topic_lines.append(
                    f"{safe_block}: выявлены затруднения в выполнении типовых заданий. Задания: {', '.join('№' + x for x in affected) if affected else 'нет'}."
                )
    else:
        topic_lines.append("Явные тематические дефициты по данным работ не выявлены.")

    weak_block_lines = []
    skill_profile = _infer_competencies(
        subject_key, list(topic_blocks.keys()), weak_candidates, second_part_avg or 0.0, first_part_avg
    )
    skill_fallback = default_skill_profile(data.subject, data.exam_type)
    for item in skill_fallback:
        if item not in skill_profile:
            skill_profile.append(item)
    skill_profile = skill_profile[:5] or skill_fallback[:3]
    for idx, (block, topics) in enumerate(topic_blocks.items()):
        inferred_skill = skill_profile[idx % len(skill_profile)]
        weak_block_lines.append(
            f"{block}: дефицит по темам ({'; '.join(topics[:2])}); "
            f"требуется усиление умения — {inferred_skill}."
        )
    if not weak_block_lines:
        weak_block_lines.append("Слабые задания не формируют выраженного блока дефицитов.")

    skill_lines = [f"Недостаточно устойчиво сформировано умение: {line}." for line in skill_profile]

    if systemic_zero_tasks or avg_task_success < 45 or data.pass_rate < 50:
        conclusion = (
            "Выявлены системные дефициты подготовки: низкая доля выполнения по ряду заданий "
            "и недостаточный уровень устойчивых предметных умений."
        )
    elif data.pass_rate >= 90 and (
        (grading_system == "5-point" and data.avg_score >= 3.8)
        or (grading_system != "5-point" and data.avg_score >= 70)
    ):
        conclusion = (
            "Общий уровень подготовки стабильный; основные предметные умения сформированы. "
            "Рекомендуется точечная работа по заданиям с относительно более низкой успешностью."
        )
    else:
        conclusion = (
            "Результаты в целом соответствуют ожидаемому уровню; "
            "для повышения устойчивости целесообразна адресная отработка отдельных тем и форматов заданий."
        )

    weak_task_nums = ", ".join(f"№{t['id']}" for t in weak_candidates[:8]) or "по данным диагностики"
    recommendations = {
        "Задания": [
            f"Провести поэтапную отработку заданий {weak_task_nums}."
        ],
        "Темы": [
            "Сформировать тематические модули по дефицитным блокам: " + ", ".join(topic_blocks.keys())
            if topic_blocks
            else "Поддерживать текущую тематическую траекторию с точечной коррекцией."
        ],
        "Навыки": [
            f"Тренировать: {skill_profile[0]}.",
            f"Развивать: {skill_profile[1]}.",
            f"Закреплять: {skill_profile[2]}.",
        ],
        "Экзаменационная стратегия": [
            "Еженедельно проводить мини-форматы экзамена с ограничением времени.",
            "Отрабатывать структуру бланка и формат ответа по критериям проверки.",
        ],
    }
    if data.exam_type == "oge":
        domain = _domain_from_subject_key(subject_key)
        if domain == "math":
            recommendations["Экзаменационная стратегия"] = [
                "Регулярные тренировки по заданиям 1–19 в формате КИМ с контролем времени.",
                "Поэтапная отработка заданий 20–25 с полным оформлением решения.",
                "Диагностика типичных вычислительных и логических ошибок.",
            ]
            recommendations["Навыки"] = list(default_skill_profile(data.subject, data.exam_type)[:3])
        elif domain == "russian":
            recommendations["Экзаменационная стратегия"] = [
                "Регулярные тренировочные задания по критериям ОГЭ в стандартном формате.",
                "Коррекционные мини-практикумы по орфографии и пунктуации.",
                "Отработка письменных форматов ответа по критериям.",
            ]
            recommendations["Навыки"] = list(default_skill_profile(data.subject, data.exam_type)[:3])
        else:
            recommendations["Экзаменационная стратегия"] = [
                "Регулярные тренировочные задания по критериям ОГЭ в стандартном формате.",
                "Коррекционные мини-практикумы на устойчивость базовых процедур.",
                "Серии коротких диагностик для стабилизации ключевых навыков.",
            ]
            recommendations["Навыки"] = list(default_skill_profile(data.subject, data.exam_type)[:3])

    control_groups = {}
    for task in weak_candidates:
        topic = task_topic_map.get(task["id"], "")
        block = _normalize_topic(topic)
        classes = _extract_classes(topic)
        classes_label = ", ".join(classes) if classes else default_grade_label(data.exam_type)
        key = (block, classes_label)
        control_groups.setdefault(key, {"tasks": [], "actions": set(), "severity": "умеренный"})
        control_groups[key]["tasks"].append(task["id"])
        sev = _task_severity(task["success_rate"])
        rank = {"умеренный": 1, "значимый": 2, "критический": 3, "рабочий": 0}
        if rank[sev] > rank[control_groups[key]["severity"]]:
            control_groups[key]["severity"] = sev
        for action in _generate_actions(block, skill_profile, sev, subject_key):
            control_groups[key]["actions"].add(action)

    control_plan = []
    for (block, classes_label), payload in sorted(control_groups.items(), key=lambda x: x[0][0]):
        control_plan.append(
            {
                "task": ", ".join(f"№{num}" for num in sorted(payload["tasks"])),
                "block": block,
                "severity": payload["severity"],
                "classes": classes_label,
                "action": "; ".join(sorted(payload["actions"])),
            }
        )

    dynamics_lines = []
    if data.exam_year == 2025 and data.dynamics:
        dynamics_lines.append("Динамика предмета за 2023–2025 гг.:")
    elif data.exam_year == 2024 and data.dynamics:
        dynamics_lines.append("Динамика предмета за 2023–2024 гг.:")
    for row in data.dynamics:
        dynamics_lines.append(
            f"{row['year']}: средний балл {row['avg_score']}, доля сдавших {row['pass_rate']}%, участников {row['students']}."
        )
    if len(data.dynamics) >= 2:
        first = data.dynamics[0]
        last = data.dynamics[-1]
        delta = round(last["avg_score"] - first["avg_score"], 2)
        trend = "рост" if delta > 0 else ("снижение" if delta < 0 else "стабильность")
        dynamics_lines.append(
            f"Итоговая динамика: {trend} среднего балла на {delta:+} пункта относительно {first['year']} года."
        )

    strongest_area = ", ".join(f"№{t['id']}" for t in strong_candidates[:3]) if strong_candidates else "не выявлены"
    weakest_area = ", ".join(f"№{t['id']}" for t in weak_candidates[:3]) if weak_candidates else "не выявлены"
    main_weak_blocks = ", ".join(list(topic_blocks.keys())[:3]) if topic_blocks else "точечные локальные дефициты"
    executive_summary = [
        f"Уровень освоения предмета оценивается как {level}; доминирует группа: {dominant_group}.",
        f"Наиболее устойчиво выполнены задания {strongest_area}; зоны наибольшего риска — {weakest_area}.",
        f"Ключевые дефицитные блоки: {main_weak_blocks}.",
        f"Приоритет вмешательства: усиление умений '{skill_profile[0]}' и '{skill_profile[1]}'.",
    ]

    readiness = _readiness_level(
        avg_score=data.avg_score,
        dominant_group=dominant_group,
        critical_count=severity_counter["критический"],
        second_part_avg=second_part_avg,
        systemic_zero_tasks=systemic_zero_tasks,
        grading_system=grading_system,
        pass_rate=data.pass_rate,
    )

    cross_skill_lines = []
    for task in weak_candidates[:8]:
        domain = _normalize_topic(task_topic_map.get(task["id"], ""))
        task_id_num = int(task["id"])
        comp = skill_profile[task_id_num % len(skill_profile)] if skill_profile else "аналитическое применение предметных знаний"
        expanded = is_expanded_answer_task(data.exam_type, subject_key, task_id_num)
        if expanded:
            cross_skill_lines.append(
                f"Снижение успешности по заданию №{task['id']} («{domain}») связано с дефицитом навыка «{comp}» "
                f"и ограничивает результат в заданиях с развёрнутым ответом."
            )
        elif _domain_from_subject_key(subject_key) == "math":
            cross_skill_lines.append(
                f"Ошибки в задании №{task['id']} («{domain}») отражают дефицит умения «{comp}» "
                f"и снижают устойчивость при переходе к более сложным форматам."
            )
        else:
            cross_skill_lines.append(
                f"Низкие результаты по заданию №{task['id']} («{domain}») указывают на недостаточную сформированность "
                f"компетенции «{comp}» в стандартном формате КИМ."
            )
    if not cross_skill_lines:
        cross_skill_lines.append("Критических межзадательных дефицитов по данным текущего среза не выявлено.")

    sections = {
        "Краткие выводы": executive_summary,
        "1. Общие результаты": [
            f"Проанализировано {data.students_count} работ.",
            f"Средний балл: {data.avg_score}; модель оценивания: {_grading_model_label_ru(grading_system)}; по шкале подготовки — {level} уровень.",
            f"Диапазон результатов: {data.min_score}–{data.max_score}. Доля преодолевших порог: {data.pass_rate}%.",
            (
                "Методическая интерпретация: результаты свидетельствуют о текущем уровне подготовки при наличии устойчивых дефицитов в заданиях повышенной сложности."
                if weak_candidates or (second_part_avg is not None and second_part_avg < 50)
                else "Методическая интерпретация: уровень подготовки в целом соответствует ожиданиям; отдельные темы требуют точечной коррекции."
            ),
        ],
        "2. Классификация обучающихся": group_lines,
        **({"3. Динамика результатов по годам": dynamics_lines} if dynamics_lines else {}),
        "4. Сильные задания и освоенные умения": strong_task_lines,
        "5. Анализ выполнения заданий": task_lines,
        "6. Тематические дефициты (по блокам)": topic_lines,
        "7. Дефициты по блокам и умениям": weak_block_lines,
        "8. Дефициты учебных умений": skill_lines,
        "8.1 Связь дефицитов заданий с формируемыми умениями": cross_skill_lines,
        "9. Ключевые проблемы": [
            "Системные: устойчиво низкая успешность по группе заданий и/или наличие заданий с 0% выполнения."
            if systemic_zero_tasks or avg_task_success < 45
            else "Системные: выраженных массовых провалов не выявлено.",
            "Локальные: точечные пробелы у части обучающихся в отдельных темах и типах заданий.",
        ],
        "10. Выводы": [
            f"{conclusion} Доминирующая группа: {dominant_group}. "
            f"Готовность к экзамену оценивается как {readiness}."
        ],
    }
    if data.exam_type == "ege" and second_part_avg is not None:
        part2_gap = round(first_part_avg - second_part_avg, 1)
        sections["5.1 Приоритет части 2 (ЕГЭ)"] = [
            f"Успешность части 1: {first_part_avg}%, части 2: {second_part_avg}%. Разрыв: {part2_gap} п.п.",
            "Ошибки части 2 связаны с дефицитом анализа условия, построения аргументации и структурирования развернутого ответа.",
            "Ограничители высокого результата: качество доказательной позиции, критериальная точность и полнота аналитического вывода.",
        ]
    elif data.exam_type == "oge" and second_part_avg is not None:
        sections["5.1 Приоритет части 2 (ОГЭ)"] = [
            f"Успешность части 1: {first_part_avg}%, части 2: {second_part_avg}%.",
            "Для ОГЭ вторая часть требует полного оформления решения и соблюдения критериев оценивания.",
        ]
    payload = {
        "sections": sections,
        "recommendations": recommendations,
        "control_plan": control_plan,
    }

    # Optional GigaChat enhancement (safe fallback if disabled/fails)
    ai_context = {
        "exam_type": data.exam_type,
        "subject": data.subject,
        "subject_key": subject_key,
        "date": data.date,
        "grading_system": grading_system,
        "students_count": data.students_count,
        "avg_score": data.avg_score,
        "min_score": data.min_score,
        "max_score": data.max_score,
        "pass_rate": data.pass_rate,
        "strong_tasks": [t["id"] for t in strong_candidates[:8]],
        "weak_tasks": [t["id"] for t in weak_candidates[:8]],
        "task_topics": {str(k): v for k, v in list(task_topic_map.items())[:30]},
        "topic_blocks": list(topic_blocks.keys())[:8],
        "skill_profile": skill_profile,
        "part1_avg": first_part_avg,
        "part2_avg": second_part_avg,
        "has_part2_data": second_part_avg is not None,
        "part2_boundary": part2_boundary,
        "systemic_zero_tasks": systemic_zero_tasks,
        "draft_executive_summary": executive_summary,
        "draft_conclusion": sections["10. Выводы"][0] if sections.get("10. Выводы") else "",
        "draft_recommendations": recommendations,
    }
    ai_result = enhance_exam_analysis(ai_context)
    if isinstance(ai_result, dict):
        exec_lines = ai_result.get("executive_summary")
        if isinstance(exec_lines, list) and exec_lines:
            payload["sections"]["Краткие выводы"] = [str(x) for x in exec_lines[:6]]
        systemic = ai_result.get("systemic_problems")
        local = ai_result.get("local_problems")
        if isinstance(systemic, list) or isinstance(local, list):
            sys_line = (
                "Системные: " + "; ".join(str(x) for x in (systemic or [])[:4])
                if isinstance(systemic, list) and systemic
                else payload["sections"]["9. Ключевые проблемы"][0]
            )
            local_line = (
                "Локальные: " + "; ".join(str(x) for x in (local or [])[:4])
                if isinstance(local, list) and local
                else payload["sections"]["9. Ключевые проблемы"][1]
            )
            payload["sections"]["9. Ключевые проблемы"] = [sys_line, local_line]
        severity = ai_result.get("severity_summary")
        if isinstance(severity, dict):
            payload["sections"]["5. Анализ выполнения заданий"].append(
                "Сводка по уровням дефицитов: "
                + f"критический — {severity.get('critical', 'нет данных')}; "
                + f"значимый — {severity.get('significant', 'нет данных')}; "
                + f"умеренный — {severity.get('moderate', 'нет данных')}."
            )
        cross_skill = ai_result.get("cross_skill_analysis")
        if isinstance(cross_skill, list) and cross_skill:
            payload["sections"]["8.1 Cross-skill интерпретация"] = [str(x).strip() for x in cross_skill if str(x).strip()][:8]
        ai_conclusion = ai_result.get("conclusion")
        if isinstance(ai_conclusion, str) and ai_conclusion.strip():
            payload["sections"]["10. Выводы"] = [ai_conclusion.strip()]
        rec_override = ai_result.get("recommendations_override")
        if isinstance(rec_override, dict):
            for key in ("Задания", "Темы", "Навыки", "Экзаменационная стратегия"):
                val = rec_override.get(key)
                if isinstance(val, list) and val:
                    payload["recommendations"][key] = [str(x) for x in val[:8]]
    if data.exam_type == "ege":
        payload["recommendations"]["Экзаменационная стратегия"].extend(
            [
                "Мастерская сочинения и развернутого ответа с поэтапной проверкой по критериям.",
                "Тренинг аналитического ответа: выделение позиции автора, тезиса и доказательной логики.",
                "Симуляции части 2 с экспертным разбором причин потери баллов.",
                "Критериальный практикум по аргументации и качеству письменного рассуждения.",
            ]
        )
        unique_strategy = []
        for line in payload["recommendations"]["Экзаменационная стратегия"]:
            if line not in unique_strategy:
                unique_strategy.append(line)
        payload["recommendations"]["Экзаменационная стратегия"] = unique_strategy[:8]
    if data.exam_type == "oge":
        # Терминология OГЭ не должна содержать упоминаний ЕГЭ.
        for title, lines in payload["sections"].items():
            payload["sections"][title] = [str(line).replace("ЕГЭ", "ОГЭ") for line in lines]
        for key, lines in payload["recommendations"].items():
            payload["recommendations"][key] = [str(line).replace("ЕГЭ", "ОГЭ") for line in lines]
        for row in payload["control_plan"]:
            row["action"] = str(row.get("action", "")).replace("ЕГЭ", "ОГЭ")

    return payload


def generate_analysis_text(data: ExamData) -> str:
    payload = _build_analysis_payload(data)
    sections = payload["sections"]
    parts = []
    for title, lines in sections.items():
        parts.append(title + "\n" + "\n".join(lines))
    rec = payload["recommendations"]
    parts.append(
        "11. Рекомендации\n"
        + "Задания: "
        + " ".join(rec["Задания"])
        + "\nТемы: "
        + " ".join(rec["Темы"])
        + "\nНавыки: "
        + " ".join(rec["Навыки"])
        + "\nЭкзаменационная стратегия: "
        + " ".join(rec["Экзаменационная стратегия"])
    )
    if payload["control_plan"]:
        parts.append("12. План контроля заданий ниже 50%")
        for row in payload["control_plan"]:
            parts.append(
                f"{row['task']} | {row['block']} | Классы: {row['classes']} | Действия: {row['action']}"
            )
    return "\n\n".join(parts)


def _docx_set_cell_shading(cell, fill_hex: str) -> None:
    from docx.oxml import OxmlElement
    from docx.oxml.ns import qn

    shd = OxmlElement("w:shd")
    shd.set(qn("w:fill"), fill_hex)
    cell._tc.get_or_add_tcPr().append(shd)


def _style_docx_table(table, header_rows: int = 1) -> None:
    """Сетка таблицы, жирный шрифт и заливка шапки — единый вид отчётов."""
    for style_name in ("Table Grid", "Light Grid Accent 1", "Сетка таблицы"):
        try:
            table.style = style_name
            break
        except (KeyError, ValueError):
            continue
    for ri in range(min(header_rows, len(table.rows))):
        for cell in table.rows[ri].cells:
            _docx_set_cell_shading(cell, "D9E2F3")
            for paragraph in cell.paragraphs:
                for run in paragraph.runs:
                    run.bold = True
                if not paragraph.runs and paragraph.text:
                    paragraph.clear()
                    paragraph.add_run(paragraph.text).bold = True


def _pdf_font_paths():
    from pathlib import Path

    return [
        p
        for p in (
            Path(r"C:\Windows\Fonts\arial.ttf"),
            Path(r"C:\Windows\Fonts\calibri.ttf"),
            Path("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf"),
            Path("/usr/share/fonts/TTF/DejaVuSans.ttf"),
        )
        if p.is_file()
    ]


def _pdf_register_cyrillic_font() -> str:
    from reportlab.pdfbase import pdfmetrics
    from reportlab.pdfbase.ttfonts import TTFont

    for path in _pdf_font_paths():
        try:
            pdfmetrics.registerFont(TTFont("AnalizGiaCyr", str(path)))
            return "AnalizGiaCyr"
        except Exception:
            continue
    return "Helvetica"


def _pdf_make_styles(font_name: str):
    from reportlab.lib.enums import TA_LEFT
    from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet

    base = getSampleStyleSheet()
    title = ParagraphStyle(
        name="AnalizTitle",
        parent=base["Heading1"],
        fontName=font_name,
        fontSize=15,
        leading=19,
        spaceAfter=10,
        alignment=TA_LEFT,
    )
    h1 = ParagraphStyle(
        name="AnalizH1",
        parent=base["Heading2"],
        fontName=font_name,
        fontSize=11,
        leading=14,
        spaceBefore=10,
        spaceAfter=6,
        alignment=TA_LEFT,
    )
    body = ParagraphStyle(
        name="AnalizBody",
        parent=base["Normal"],
        fontName=font_name,
        fontSize=9,
        leading=12,
        alignment=TA_LEFT,
    )
    small = ParagraphStyle(
        name="AnalizSmall",
        parent=body,
        fontSize=8,
        leading=11,
    )
    return title, h1, body, small


def _pdf_p(text: str, style):
    from xml.sax.saxutils import escape

    from reportlab.platypus import Paragraph

    safe = escape(str(text or "")).replace("\n", "<br/>")
    return Paragraph(safe, style)


def _pdf_table(headers: list[str], rows: list[list], col_widths: list[float], font_name: str):
    from reportlab.lib import colors
    from reportlab.platypus import Table, TableStyle

    data = [headers] + rows
    tbl = Table(data, colWidths=col_widths, repeatRows=1)
    tbl.setStyle(
        TableStyle(
            [
                ("FONT", (0, 0), (-1, -1), font_name, 9),
                ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#D9E2F3")),
                ("TEXTCOLOR", (0, 0), (-1, 0), colors.HexColor("#0f172a")),
                ("FONT", (0, 0), (-1, 0), font_name, 9),
                ("ALIGN", (0, 0), (-1, -1), "LEFT"),
                ("VALIGN", (0, 0), (-1, -1), "MIDDLE"),
                ("GRID", (0, 0), (-1, -1), 0.4, colors.HexColor("#94a3b8")),
                ("ROWBACKGROUNDS", (0, 1), (-1, -1), [colors.white, colors.HexColor("#f8fafc")]),
                ("TOPPADDING", (0, 0), (-1, -1), 4),
                ("BOTTOMPADDING", (0, 0), (-1, -1), 4),
                ("LEFTPADDING", (0, 0), (-1, -1), 5),
                ("RIGHTPADDING", (0, 0), (-1, -1), 5),
            ]
        )
    )
    return tbl


def _pdf_build_document(story: list) -> BytesIO:
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.units import cm
    from reportlab.platypus import SimpleDocTemplate

    buf = BytesIO()
    doc = SimpleDocTemplate(
        buf,
        pagesize=A4,
        leftMargin=1.7 * cm,
        rightMargin=1.7 * cm,
        topMargin=1.4 * cm,
        bottomMargin=1.4 * cm,
    )
    doc.build(story)
    buf.seek(0)
    return buf


def generate_word_doc(data: ExamData) -> BytesIO:
    from users.report_ui.school_exam_analysis_docx import render_exam_analysis_docx

    return render_exam_analysis_docx(data)


def generate_presentation(data: ExamData) -> BytesIO:
    from pptx import Presentation
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE

    prs = Presentation()

    def add_title_slide(title: str, subtitle: str):
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = title
        slide.placeholders[1].text = subtitle

    def add_bullets(title: str, lines: list[str]):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = title
        tf = slide.placeholders[1].text_frame
        tf.clear()
        for idx, line in enumerate(lines):
            p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
            p.text = line

    def add_task_success_chart():
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = "Успешность по заданиям"
        chart_data = CategoryChartData()
        labels = [f"№{t['id']}" for t in data.tasks]
        chart_data.categories = labels
        chart_data.add_series("Успешность, %", [t["success_rate"] for t in data.tasks])
        slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED,
            x=500000,
            y=1400000,
            cx=8500000,
            cy=4200000,
            chart_data=chart_data,
        )

    def add_distribution_chart():
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = "Распределение обучающихся по уровням"
        subject_key = _subject_key(data.subject)
        groups = _score_groups(data.score_values, subject_key)
        chart_data = CategoryChartData()
        if subject_key == "math_basic":
            chart_data.categories = ["2", "3", "4", "5"]
            chart_data.add_series(
                "Количество",
                [groups["grade_2"], groups["grade_3"], groups["grade_4"], groups["grade_5"]],
            )
        else:
            chart_data.categories = ["Слабый", "Базовый", "Хороший", "Высокий"]
            chart_data.add_series("Количество", [groups["weak"], groups["basic"], groups["good"], groups["high"]])
        slide.shapes.add_chart(
            XL_CHART_TYPE.BAR_CLUSTERED,
            x=500000,
            y=1400000,
            cx=8500000,
            cy=4200000,
            chart_data=chart_data,
        )

    def add_heatmap_like_slide():
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = "Карта рисков по заданиям (heatmap-представление)"
        rows = min(len(data.tasks) + 1, 16)
        cols = 3
        table = slide.shapes.add_table(rows, cols, 500000, 1400000, 8500000, 4200000).table
        table.cell(0, 0).text = "Задание"
        table.cell(0, 1).text = "Успешность"
        table.cell(0, 2).text = "Зона"
        for idx, task in enumerate(data.tasks[:15], start=1):
            rate = task["success_rate"]
            zone = "Высокая" if rate >= 80 else ("Рабочая" if rate >= 50 else "Риск")
            table.cell(idx, 0).text = f"№{task['id']}"
            table.cell(idx, 1).text = f"{rate}%"
            table.cell(idx, 2).text = zone

    payload = _build_analysis_payload(data)
    sections = payload["sections"]
    recommendations = payload["recommendations"]
    control_plan = payload["control_plan"]

    add_title_slide(f"Анализ результатов {_exam_type_label_ru(data.exam_type)} по {data.subject}", f"Дата экзамена: {data.date}")
    add_bullets(
        "Общие результаты",
        [
            f"Участников: {data.students_count}",
            f"Средний балл: {data.avg_score}",
            f"Минимум/максимум: {data.min_score}/{data.max_score}",
            f"Доля сдавших: {data.pass_rate}%",
        ],
    )
    add_bullets("Распределение результатов", sections.get("2. Классификация обучающихся", []))
    add_distribution_chart()
    if sections.get("3. Динамика результатов по годам"):
        add_bullets("Динамика по годам", sections.get("3. Динамика результатов по годам", []))
    add_bullets("Сильные задания", sections.get("4. Сильные задания и освоенные умения", []))
    add_bullets("Анализ заданий", sections.get("5. Анализ выполнения заданий", []))
    add_task_success_chart()
    add_heatmap_like_slide()
    add_bullets("Проблемные зоны", sections.get("9. Ключевые проблемы", []))
    add_bullets("Тематические дефициты", sections.get("6. Тематические дефициты (по блокам)", []))
    add_bullets("Дефициты по умениям", sections.get("8. Дефициты учебных умений", []))
    add_bullets("Выводы", sections.get("10. Выводы", []))
    add_bullets(
        "Рекомендации",
        recommendations.get("Задания", [])
        + recommendations.get("Темы", [])
        + recommendations.get("Навыки", [])
        + recommendations.get("Экзаменационная стратегия", []),
    )
    if control_plan:
        control_lines = [
            f"{item['task']} | {item['block']} | {item.get('severity', 'значимый')} | {item['classes']} кл. | {item['action']}"
            for item in control_plan[:8]
        ]
    else:
        control_lines = ["Заданий ниже 50% не выявлено."]
    add_bullets("План контроля заданий <50%", control_lines)

    output = BytesIO()
    prs.save(output)
    output.seek(0)
    return output


def generate_xlsx_report(data: ExamData) -> BytesIO:
    from openpyxl import Workbook
    from openpyxl.styles import Font

    payload = _build_analysis_payload(data)
    wb = Workbook()
    ws = wb.active
    ws.title = "Аналитика"

    row = 1
    ws.cell(row=row, column=1, value="АНАЛИТИЧЕСКАЯ СПРАВКА").font = Font(bold=True, size=14)
    row += 1
    ws.cell(row=row, column=1, value=f"{_exam_type_label_ru(data.exam_type)} | {data.subject} | {data.date}")
    row += 2

    ws.cell(row=row, column=1, value="Сводные показатели").font = Font(bold=True)
    row += 1
    ws.append(["Участников", "Средний балл", "Мин", "Макс", "Доля сдавших"])
    ws.append([data.students_count, data.avg_score, data.min_score, data.max_score, f"{data.pass_rate}%"])
    row = ws.max_row + 2

    ws.cell(row=row, column=1, value="Анализ выполнения заданий").font = Font(bold=True)
    row += 1
    ws.append(["Задание", "Успешность, %", "Верно (+)", "Ошибок (-)", "Всего"])
    for task in data.tasks:
        ws.append([f"№{task['id']}", task["success_rate"], task["correct"], task["wrong"], task["total"]])
    row = ws.max_row + 2

    ws.cell(row=row, column=1, value="Рекомендации").font = Font(bold=True)
    row += 1
    ws.append(["Блок", "Текст"])
    recommendations = payload.get("recommendations", {})
    for block_name, lines in recommendations.items():
        for line in lines:
            ws.append([block_name, line])
    row = ws.max_row + 2

    ws.cell(row=row, column=1, value="План контроля заданий ниже 50%").font = Font(bold=True)
    row += 1
    ws.append(["Задание", "Тематический блок", "Уровень дефицита", "Классы", "Действия"])
    control_rows = payload.get("control_plan", [])
    if control_rows:
        for item in control_rows:
            ws.append(
                [
                    item.get("task", ""),
                    item.get("block", ""),
                    item.get("severity", ""),
                    item.get("classes", ""),
                    item.get("action", ""),
                ]
            )
    else:
        ws.append(["-", "Заданий с успешностью ниже 50% не выявлено.", "", "", ""])

    for col in ("A", "B", "C", "D", "E", "F"):
        ws.column_dimensions[col].width = 28

    output = BytesIO()
    wb.save(output)
    output.seek(0)
    return output


def generate_pdf_report(data: ExamData) -> BytesIO:
    from reportlab.lib.units import mm
    from reportlab.platypus import Spacer

    font_name = _pdf_register_cyrillic_font()
    title_s, h1_s, body_s, _small = _pdf_make_styles(font_name)
    payload = _build_analysis_payload(data)
    story = []
    story.append(_pdf_p("АНАЛИТИЧЕСКАЯ СПРАВКА", title_s))
    story.append(_pdf_p(f"{_exam_type_label_ru(data.exam_type)} | {data.subject} | {data.date}", body_s))
    story.append(Spacer(1, 4))

    summary_rows = [
        [
            str(data.students_count),
            str(data.avg_score),
            str(data.min_score),
            str(data.max_score),
            f"{data.pass_rate}%",
        ]
    ]
    story.append(_pdf_p("Сводные показатели", h1_s))
    story.append(
        _pdf_table(
            ["Участников", "Средний", "Мин", "Макс", "Доля сдавших"],
            summary_rows,
            [28 * mm, 28 * mm, 22 * mm, 22 * mm, 32 * mm],
            font_name,
        )
    )
    story.append(Spacer(1, 8))

    if data.tasks:
        story.append(_pdf_p("Успешность по заданиям", h1_s))
        task_rows = [
            [f"№{t['id']}", str(t.get("success_rate", "")), str(t.get("correct", "")), str(t.get("wrong", "")), str(t.get("total", ""))]
            for t in data.tasks
        ]
        story.append(
            _pdf_table(
                ["Задание", "Успешность, %", "Верно (+)", "Ошибок (−)", "Всего"],
                task_rows,
                [22 * mm, 30 * mm, 28 * mm, 28 * mm, 28 * mm],
                font_name,
            )
        )
        story.append(Spacer(1, 8))

    story.append(_pdf_p("Ключевые выводы", h1_s))
    for line in payload.get("sections", {}).get("Краткие выводы", [])[:10]:
        story.append(_pdf_p(f"• {line}", body_s))
    story.append(Spacer(1, 6))

    story.append(_pdf_p("Рекомендации", h1_s))
    for block, lines in payload.get("recommendations", {}).items():
        story.append(_pdf_p(f"<b>{block}</b>", body_s))
        for line in lines[:6]:
            story.append(_pdf_p(f"  — {line}", body_s))
    story.append(Spacer(1, 6))

    story.append(_pdf_p("План контроля заданий ниже 50%", h1_s))
    control_rows = payload.get("control_plan", [])
    if control_rows:
        cr = []
        for item in control_rows:
            cr.append(
                [
                    str(item.get("task", ""))[:20],
                    str(item.get("block", ""))[:42],
                    str(item.get("severity", ""))[:14],
                    str(item.get("classes", ""))[:12],
                    str(item.get("action", ""))[:52],
                ]
            )
        story.append(
            _pdf_table(
                ["Задание", "Тематический блок", "Дефицит", "Классы", "Действия"],
                cr,
                [22 * mm, 52 * mm, 24 * mm, 22 * mm, 52 * mm],
                font_name,
            )
        )
    else:
        story.append(_pdf_p("Заданий с успешностью ниже 50% не выявлено.", body_s))

    return _pdf_build_document(story)


def generate_school_gia_summary_xlsx(school_id: int) -> BytesIO:
    from openpyxl import Workbook
    from openpyxl.styles import Font

    results = ExamResult.objects.filter(student__school_id=school_id).select_related("exam", "student")

    wb = Workbook()
    ws = wb.active
    ws.title = "Свод ГИА"

    ws["A1"] = "СВОД РЕЗУЛЬТАТОВ ГИА (ОО)"
    ws["A1"].font = Font(bold=True, size=14)

    if not results.exists():
        ws["A3"] = "Недостаточно данных для формирования отчета."
        output = BytesIO()
        wb.save(output)
        output.seek(0)
        return output

    total = results.count()
    avg_score = results.aggregate(v=Avg("score"))["v"] or 0
    pass_count = results.filter(passed=True).count()
    pass_rate = round((pass_count / total) * 100, 1) if total else 0.0

    # Quality indicator: share of strong results.
    max_score = results.aggregate(v=Max("score"))["v"] or 100
    quality_threshold = 4 if float(max_score) <= 5 else 60
    quality_count = results.filter(score__gte=quality_threshold).count()
    quality_rate = round((quality_count / total) * 100, 1) if total else 0.0

    ws["A3"] = "Краткая статистика"
    ws["A3"].font = Font(bold=True)
    ws.append(["Показатель", "Значение"])
    ws.append(["Число участников", total])
    ws.append(["Средний балл", round(float(avg_score), 2)])
    ws.append(["Качество знаний, %", quality_rate])
    ws.append(["Успеваемость, %", pass_rate])

    # Distribution of grades/scores.
    ws.append([])
    ws.append(["Распределение оценок / уровней", "Количество"])
    max_score_value = float(max_score or 0)
    if max_score_value <= 5:
        ws.append(["Оценка 2", results.filter(score__lt=3).count()])
        ws.append(["Оценка 3", results.filter(score__gte=3, score__lt=4).count()])
        ws.append(["Оценка 4", results.filter(score__gte=4, score__lt=5).count()])
        ws.append(["Оценка 5", results.filter(score__gte=5).count()])
    else:
        ws.append(["0-35", results.filter(score__lte=35).count()])
        ws.append(["36-60", results.filter(score__gt=35, score__lte=60).count()])
        ws.append(["61-80", results.filter(score__gt=60, score__lte=80).count()])
        ws.append(["81-100", results.filter(score__gt=80).count()])

    # Subject results.
    ws.append([])
    ws.append(["Результаты по предметам", "Участников", "Средний балл", "Успеваемость, %"])
    subject_rows = (
        results.values("exam__subject")
        .annotate(cnt=Count("id"), avg=Avg("score"), passed=Count("id", filter=Q(passed=True)))
        .order_by("exam__subject")
    )
    for row in subject_rows:
        cnt = int(row["cnt"] or 0)
        passed = int(row["passed"] or 0)
        ws.append(
            [
                row["exam__subject"] or "Предмет не указан",
                cnt,
                round(float(row["avg"] or 0), 2),
                round((passed / cnt) * 100, 1) if cnt else 0.0,
            ]
        )

    # Class ranking (by average score).
    ws.append([])
    ws.append(["Рейтинг классов", "Участников", "Средний балл", "Успеваемость, %"])
    class_rows = (
        results.values("student__grade")
        .annotate(cnt=Count("id"), avg=Avg("score"), passed=Count("id", filter=Q(passed=True)))
        .order_by("-avg", "student__grade")
    )
    for row in class_rows:
        cnt = int(row["cnt"] or 0)
        passed = int(row["passed"] or 0)
        ws.append(
            [
                (row["student__grade"] or "Класс не указан"),
                cnt,
                round(float(row["avg"] or 0), 2),
                round((passed / cnt) * 100, 1) if cnt else 0.0,
            ]
        )

    ws.column_dimensions["A"].width = 38
    ws.column_dimensions["B"].width = 18
    ws.column_dimensions["C"].width = 18
    ws.column_dimensions["D"].width = 20

    output = BytesIO()
    wb.save(output)
    output.seek(0)
    return output


def generate_school_gia_summary_docx(school_id: int, exam_type: str, year: int | None = None) -> BytesIO:
    from docx import Document

    def _is_passed_for_summary(row: dict, threshold_cache: dict) -> bool:
        from exams.passing import ege_result_passed, oge_score_passed

        if et != "ege":
            return oge_score_passed(row.get("score"), row.get("passed"))
        return ege_result_passed(
            subject_name=row.get("exam__subject"),
            year=row.get("exam__year"),
            score=row.get("score"),
            passed_flag=row.get("passed"),
            exam_code=row.get("exam__code"),
            cache=threshold_cache,
        )

    et = (exam_type or "").strip().lower()
    if et not in {"ege", "oge"}:
        et = "ege"

    results = ExamResult.objects.filter(student__school_id=school_id, exam__exam_type=et).select_related("exam", "student")
    if year:
        results = results.filter(exam__year=year)

    from users.report_ui.school_gia_summary import build_gia_summary_presentation

    doc = Document()
    title_label = "ЕГЭ" if et == "ege" else "ОГЭ"

    if not results.exists():
        doc.add_paragraph("Недостаточно данных для формирования отчета.")
        output = BytesIO()
        doc.save(output)
        output.seek(0)
        return output

    results_rows = list(
        results.values(
            "exam__subject",
            "exam__code",
            "exam__year",
            "student__grade",
            "score",
            "passed",
        )
    )
    threshold_cache = {}
    pass_count = 0
    passed_by_subject = {}
    totals_by_subject = {}
    passed_by_class = {}
    totals_by_class = {}
    for row in results_rows:
        is_passed = _is_passed_for_summary(row, threshold_cache)
        subject_name = row.get("exam__subject") or "Предмет не указан"
        class_name = row.get("student__grade") or "Класс не указан"
        totals_by_subject[subject_name] = totals_by_subject.get(subject_name, 0) + 1
        totals_by_class[class_name] = totals_by_class.get(class_name, 0) + 1
        if is_passed:
            pass_count += 1
            passed_by_subject[subject_name] = passed_by_subject.get(subject_name, 0) + 1
            passed_by_class[class_name] = passed_by_class.get(class_name, 0) + 1

    total = results.count()
    participants_unique = results.values("student_id").distinct().count()
    avg_score = float(results.aggregate(v=Avg("score"))["v"] or 0)
    pass_rate = round((pass_count / total) * 100, 1) if total else 0.0

    # Hard split by exam type: EGE uses score model, OGE uses grade model.
    quality_threshold = 4 if et == "oge" else 60
    high_threshold = 5 if et == "oge" else 70
    quality_count = results.filter(score__gte=quality_threshold).count()
    quality_rate = round((quality_count / total) * 100, 1) if total else 0.0
    high_count = results.filter(score__gte=high_threshold).count()
    failed_count = total - pass_count

    if et == "oge":
        distribution = [
            {"label": "2", "value": results.filter(score__lt=3).count()},
            {"label": "3", "value": results.filter(score__gte=3, score__lt=4).count()},
            {"label": "4", "value": results.filter(score__gte=4, score__lt=5).count()},
            {"label": "5", "value": results.filter(score__gte=5).count()},
        ]
    else:
        distribution = [
            {"label": "0-35", "value": results.filter(score__lte=35).count()},
            {"label": "36-60", "value": results.filter(score__gt=35, score__lte=60).count()},
            {"label": "61-80", "value": results.filter(score__gt=60, score__lte=80).count()},
            {"label": "81-100", "value": results.filter(score__gt=80).count()},
        ]

    subject_qs = (
        results.values("exam__subject")
        .annotate(cnt=Count("id"), avg=Avg("score"))
        .order_by("exam__subject")
    )
    subject_rows = []
    for row_data in subject_qs:
        cnt = int(row_data["cnt"] or 0)
        subject_name = row_data["exam__subject"] or "Предмет не указан"
        passed = int(passed_by_subject.get(subject_name, 0))
        subject_rows.append(
            {
                "exam__subject": subject_name,
                "participants": cnt,
                "avg": round(float(row_data["avg"] or 0), 2),
                "pass_rate": round((passed / cnt) * 100, 1) if cnt else 0.0,
            }
        )

    avg_label = "Средняя оценка" if et == "oge" else "Средний балл"
    ui = build_gia_summary_presentation(
        exam_type=et,
        year=year,
        kpis={
            "participants": participants_unique,
            "total_results": total,
            "avg_score": round(avg_score, 2),
            "avg_label": avg_label,
            "quality_rate": quality_rate,
            "pass_rate": pass_rate,
            "high_count": high_count,
            "failed_count": failed_count,
            "risk_students": failed_count,
        },
        distribution=distribution,
        subject_rows=subject_rows,
    )

    class_rows_ui = []
    class_qs = (
        results.values("student__grade")
        .annotate(cnt=Count("id"), avg=Avg("score"))
        .order_by("-avg", "student__grade")
    )
    for row_data in class_qs:
        cnt = int(row_data["cnt"] or 0)
        class_name = row_data["student__grade"] or "Класс не указан"
        passed = int(passed_by_class.get(class_name, 0))
        class_rows_ui.append(
            {
                "name": class_name,
                "participants": cnt,
                "avg": round(float(row_data["avg"] or 0), 2),
                "pass_rate": round((passed / cnt) * 100, 1) if cnt else 0.0,
            }
        )

    from users.report_ui.school_gia_summary_docx import render_gia_summary_docx

    render_gia_summary_docx(
        doc,
        ui,
        exam_label=title_label,
        year=year,
        avg_label=avg_label,
        class_rows=class_rows_ui,
    )

    output = BytesIO()
    doc.save(output)
    output.seek(0)
    return output


def _build_school_info_stat_payload(school_id: int, exam_type: str, year: int | None = None) -> dict:
    et = (exam_type or "").strip().lower()
    if et not in {"ege", "oge"}:
        et = "ege"

    qs = ExamResult.objects.filter(student__school_id=school_id, exam__exam_type=et)
    if year:
        qs = qs.filter(exam__year=year)

    if not qs.exists():
        return {"has_data": False, "message": "Недостаточно данных для формирования отчета."}

    def _threshold_subject_key(subject_name: str) -> str | None:
        title = (subject_name or "").strip().lower()
        if "рус" in title:
            return "russian"
        if "математика" in title and "проф" in title:
            return "math_profile"
        if "математика" in title and "баз" in title:
            return "math_basic"
        if "обществ" in title:
            return "social"
        if "информат" in title:
            return "informatics"
        if "физик" in title:
            return "physics"
        if "хими" in title:
            return "chemistry"
        if "биолог" in title:
            return "biology"
        if "истори" in title:
            return "history"
        if "литератур" in title:
            return "literature"
        if "географ" in title:
            return "geography"
        if any(lang in title for lang in ("англий", "немец", "француз", "испан", "китай", "иностран")):
            return "foreign_language"
        return None

    threshold_cache: dict[tuple[int, str], EgePassingThreshold | None] = {}

    def _is_passed_ege(row: dict) -> bool:
        if is_gve_exam(exam_code=row.get("exam__code"), subject_name=row.get("exam__subject")):
            score_value = float(row.get("score") or 0)
            return score_value >= 3
        subject_key = _threshold_subject_key(row.get("exam__subject") or "")
        if not subject_key:
            return bool(row.get("passed"))
        threshold_key = (int(row.get("exam__year") or 0), subject_key)
        if threshold_key not in threshold_cache:
            threshold_cache[threshold_key] = (
                EgePassingThreshold.objects.filter(year=threshold_key[0], subject_key=subject_key)
                .only("minimum_score", "minimum_grade")
                .first()
            )
        threshold = threshold_cache[threshold_key]
        if not threshold:
            return bool(row.get("passed"))
        score_value = float(row.get("score") or 0)
        if threshold.minimum_score is not None:
            return score_value >= float(threshold.minimum_score)
        if threshold.minimum_grade is not None:
            return score_value >= float(threshold.minimum_grade)
        return bool(row.get("passed"))

    total = qs.count()
    participants = qs.values("student_id").distinct().count()
    ege_rows = []
    if et == "ege":
        ege_rows = list(
            qs.values(
                "student_id",
                "student__grade",
                "exam__subject",
                "exam__code",
                "exam__year",
                "score",
                "passed",
            )
        )

    # Средний/мин/макс по стобалльной шкале ЕГЭ — без оценок ГВЭ (2–5), иначе min=3.
    metric_qs = qs
    if et == "ege" and ege_rows:
        gve_ids = [
            i
            for i, row in enumerate(ege_rows)
            if is_gve_exam(exam_code=row.get("exam__code"), subject_name=row.get("exam__subject"))
        ]
        if gve_ids and len(gve_ids) < len(ege_rows):
            metric_qs = qs.exclude(
                Q(exam__code__in=["51", "051"]) | Q(exam__subject__icontains="ГВЭ") | Q(exam__subject__icontains="гвэ")
            )

    avg_score = float(metric_qs.aggregate(v=Avg("score"))["v"] or 0)
    min_score = float(metric_qs.aggregate(v=Min("score"))["v"] or 0)
    max_score = float(metric_qs.aggregate(v=Max("score"))["v"] or 0)
    passed = qs.filter(passed=True).count()
    if et == "ege":
        passed = sum(1 for row in ege_rows if _is_passed_ege(row))
    pass_rate = round((passed / total) * 100, 1) if total else 0.0
    quality_threshold = 4 if max_score <= 5 else 60
    high_threshold = 5 if max_score <= 5 else 70
    if et == "ege":
        quality_count = 0
        high_count = 0
        for row in ege_rows:
            score = float(row.get("score") or 0)
            if is_gve_exam(exam_code=row.get("exam__code"), subject_name=row.get("exam__subject")):
                if score >= 4:
                    quality_count += 1
                if score >= 5:
                    high_count += 1
            else:
                if score >= quality_threshold:
                    quality_count += 1
                if score >= high_threshold:
                    high_count += 1
    else:
        quality_count = qs.filter(score__gte=quality_threshold).count()
        high_count = qs.filter(score__gte=high_threshold).count()
    quality_rate = round((quality_count / total) * 100, 1) if total else 0.0
    failed_count = max(total - passed, 0)

    prev_avg = None
    prev_pass_rate = None
    if year:
        prev_qs = ExamResult.objects.filter(student__school_id=school_id, exam__exam_type=et, exam__year=year - 1)
        prev_metric = prev_qs
        if et == "ege":
            prev_metric = prev_qs.exclude(
                Q(exam__code__in=["51", "051"]) | Q(exam__subject__icontains="ГВЭ") | Q(exam__subject__icontains="гвэ")
            )
            if not prev_metric.exists():
                prev_metric = prev_qs
        prev = prev_metric.aggregate(v=Avg("score"))["v"]
        if prev is not None:
            prev_avg = round(float(prev), 2)
        prev_total = prev_qs.count()
        if prev_total:
            if et == "ege":
                prev_rows = list(prev_qs.values("exam__subject", "exam__code", "exam__year", "score", "passed"))
                prev_passed = sum(1 for row in prev_rows if _is_passed_ege(row))
            else:
                prev_passed = prev_qs.filter(passed=True).count()
            prev_pass_rate = round((prev_passed / prev_total) * 100, 1)

    subject_rows = []
    if et == "ege":
        by_subject: dict[str, dict] = {}
        for row in ege_rows:
            subject_name = gve_subject_label(row.get("exam__subject"), row.get("exam__code"))
            is_gve = is_gve_exam(exam_code=row.get("exam__code"), subject_name=row.get("exam__subject"))
            bucket = by_subject.setdefault(
                subject_name,
                {
                    "exam__subject": subject_name,
                    "student_ids": set(),
                    "results": 0,
                    "sum_score": 0.0,
                    "min_v": None,
                    "max_v": None,
                    "passed": 0,
                    "quality": 0,
                    "is_gve": is_gve,
                },
            )
            score = float(row.get("score") or 0)
            sid = row.get("student_id")
            if sid is not None:
                bucket["student_ids"].add(sid)
            bucket["results"] += 1
            bucket["sum_score"] += score
            bucket["min_v"] = score if bucket["min_v"] is None else min(bucket["min_v"], score)
            bucket["max_v"] = score if bucket["max_v"] is None else max(bucket["max_v"], score)
            if _is_passed_ege(row):
                bucket["passed"] += 1
            q_thr = 4 if is_gve or max_score <= 5 else quality_threshold
            if score >= q_thr:
                bucket["quality"] += 1
        subject_rows = sorted(by_subject.values(), key=lambda x: x["exam__subject"])
        for row in subject_rows:
            results_n = int(row.pop("results", 0) or 0)
            student_ids = row.pop("student_ids", set()) or set()
            row.pop("is_gve", None)
            participants_n = len(student_ids) or results_n
            row["participants"] = participants_n
            row["avg"] = (row["sum_score"] / results_n) if results_n else 0.0
            row["quality_rate"] = round((int(row["quality"] or 0) / results_n) * 100, 1) if results_n else 0.0
            row["pass_rate"] = round((int(row["passed"] or 0) / results_n) * 100, 1) if results_n else 0.0
            row["failed"] = max(results_n - int(row["passed"] or 0), 0)
            del row["sum_score"]
            del row["quality"]
            del row["passed"]
    else:
        subject_rows = list(
            qs.values("exam__subject")
            .annotate(
                participants=Count("student_id", distinct=True),
                results=Count("id"),
                avg=Avg("score"),
                min_v=Min("score"),
                max_v=Max("score"),
                passed=Count("id", filter=Q(passed=True)),
                quality=Count("id", filter=Q(score__gte=quality_threshold)),
            )
            .order_by("exam__subject")
        )
        for row in subject_rows:
            results_n = int(row.pop("results", 0) or 0) or int(row.get("participants") or 0)
            row["quality_rate"] = round((int(row["quality"] or 0) / results_n) * 100, 1) if results_n else 0.0
            row["pass_rate"] = round((int(row["passed"] or 0) / results_n) * 100, 1) if results_n else 0.0
            row["failed"] = max(results_n - int(row["passed"] or 0), 0)

    subjects_count = len(subject_rows)

    # Score distribution bins (без ГВЭ на стобалльной шкале ЕГЭ)
    dist_qs = metric_qs if et == "ege" else qs
    if max_score <= 5:
        bins = [("2", dist_qs.filter(score__lt=3).count()), ("3", dist_qs.filter(score__gte=3, score__lt=4).count()),
                ("4", dist_qs.filter(score__gte=4, score__lt=5).count()), ("5", dist_qs.filter(score__gte=5).count())]
    else:
        bins = [
            ("90-100", dist_qs.filter(score__gte=90).count()),
            ("80-89", dist_qs.filter(score__gte=80, score__lt=90).count()),
            ("70-79", dist_qs.filter(score__gte=70, score__lt=80).count()),
            ("60-69", dist_qs.filter(score__gte=60, score__lt=70).count()),
            ("50-59", dist_qs.filter(score__gte=50, score__lt=60).count()),
            ("0-49", dist_qs.filter(score__lt=50).count()),
        ]

    class_rows = []
    if et == "ege":
        by_class: dict[str, dict] = {}
        for row in ege_rows:
            class_name = row.get("student__grade") or "Класс не указан"
            bucket = by_class.setdefault(
                class_name,
                {"student__grade": class_name, "student_ids": set(), "results": 0, "sum_score": 0.0, "passed": 0},
            )
            score = float(row.get("score") or 0)
            sid = row.get("student_id")
            if sid is not None:
                bucket["student_ids"].add(sid)
            bucket["results"] += 1
            bucket["sum_score"] += score
            if _is_passed_ege(row):
                bucket["passed"] += 1
        class_rows = list(by_class.values())
        for row in class_rows:
            results_n = int(row.pop("results", 0) or 0)
            student_ids = row.pop("student_ids", set()) or set()
            participants_n = len(student_ids) or results_n
            row["participants"] = participants_n
            row["avg"] = (row["sum_score"] / results_n) if results_n else 0.0
            row["pass_rate"] = round((int(row["passed"] or 0) / results_n) * 100, 1) if results_n else 0.0
            del row["sum_score"]
            del row["passed"]
        class_rows.sort(key=lambda x: (-float(x.get("avg") or 0), str(x.get("student__grade") or "")))
    else:
        class_rows = list(
            qs.values("student__grade")
            .annotate(
                participants=Count("student_id", distinct=True),
                results=Count("id"),
                avg=Avg("score"),
                passed=Count("id", filter=Q(passed=True)),
            )
            .order_by("-avg", "student__grade")
        )
        for row in class_rows:
            results_n = int(row.pop("results", 0) or 0) or int(row.get("participants") or 0)
            row["pass_rate"] = round((int(row["passed"] or 0) / results_n) * 100, 1) if results_n else 0.0

    # Comparisons (school vs district vs republic)
    district_avg = None
    republic_avg = None
    school = School.objects.filter(id=school_id).only("district_id").first()
    if school and school.district_id:
        district_qs = ExamResult.objects.filter(
            student__school__district_id=school.district_id,
            exam__exam_type=et,
            exam__year=year if year else qs.first().exam.year,
        )
        if district_qs.exists():
            district_avg = round(float(district_qs.aggregate(v=Avg("score"))["v"] or 0), 2)
    republic_qs = ExamResult.objects.filter(
        exam__exam_type=et,
        exam__year=year if year else qs.first().exam.year,
    )
    if republic_qs.exists():
        republic_avg = round(float(republic_qs.aggregate(v=Avg("score"))["v"] or 0), 2)

    # OGE task analytics for production dashboard
    task_rows = []
    if et == "oge":
        task_qs = TaskResult.objects.filter(student__school_id=school_id, exam__exam_type="oge")
        if year:
            task_qs = task_qs.filter(exam__year=year)
        task_rows = list(
            task_qs.values("exam__subject", "task_number")
            .annotate(total=Count("id"), plus=Count("id", filter=~Q(value__in=["-", "0", ""])))
            .order_by("exam__subject", "task_number")
        )
        for row in task_rows:
            total_task = int(row["total"] or 0)
            plus_task = int(row["plus"] or 0)
            row["success_rate"] = round((plus_task / total_task) * 100, 1) if total_task else 0.0
            row["risk"] = _risk_level_ru_from_rates(float(row["success_rate"] or 0))

    if et == "ege":
        risk_students = {row.get("student_id") for row in ege_rows if not _is_passed_ege(row)}
        risk_count = len([s for s in risk_students if s is not None])
    else:
        risk_count = qs.filter(Q(passed=False) | Q(score__lt=3)).values("student_id").distinct().count()
    weak_subjects = sorted(
        [r for r in subject_rows if _is_weak_subject_row(r, exam_type=et, max_score=max_score)],
        key=lambda x: (x["pass_rate"], float(x["avg"] or 0)),
    )
    weak_tasks = sorted(task_rows, key=lambda x: x["success_rate"])[:12] if task_rows else []
    avg_delta = round(avg_score - prev_avg, 2) if prev_avg is not None else None
    pass_delta = round(pass_rate - prev_pass_rate, 1) if prev_pass_rate is not None else None

    ai_insights = []
    if avg_delta is not None:
        ai_insights.append(
            f"Динамика среднего результата к прошлому году: {'+' if avg_delta > 0 else ''}{avg_delta}."
        )
    if weak_subjects:
        ai_insights.append(
            "Проблемные предметы: " + ", ".join((r["exam__subject"] or "предмет") for r in weak_subjects[:3]) + "."
        )
    if weak_tasks:
        ai_insights.append(
            "Критические задания ОГЭ: " + ", ".join(f"{r['exam__subject']} №{r['task_number']}" for r in weak_tasks[:5]) + "."
        )
    if risk_count:
        ai_insights.append(f"Группа риска: {risk_count} обучающихся.")

    recommendations = []
    if weak_subjects:
        recommendations.append(
            "Адресная коррекция по предметам: "
            + ", ".join(
                f"{(r.get('exam__subject') or 'предмет')} (усп. {r.get('pass_rate')}%)"
                for r in weak_subjects[:4]
            )
            + "."
        )
    if weak_tasks:
        recommendations.append(
            "Повторная диагностика по заданиям: "
            + ", ".join(f"{r['exam__subject']} №{r['task_number']} ({r['success_rate']}%)" for r in weak_tasks[:5])
            + "."
        )
    if risk_count:
        recommendations.append(f"Индивидуальные маршруты для группы риска: {risk_count} обучающихся.")
    if avg_delta is not None and avg_delta < 0:
        recommendations.append(f"Остановить снижение среднего балла к прошлому году ({avg_delta:+}).")
    if not recommendations:
        recommendations.append(
            f"Поддерживать текущий уровень: средний {round(avg_score, 2)}, успеваемость {pass_rate}%."
        )
    # Без GigaChat: insights и рекомендации только из фактических метрик БД.

    dynamics = []
    if year:
        years = [year - 2, year - 1, year]
        if et == "ege":
            dyn_rows = list(
                ExamResult.objects.filter(student__school_id=school_id, exam__exam_type=et, exam__year__in=years)
                .values("exam__year", "exam__subject", "exam__code", "student_id", "score", "passed")
            )
            by_year = {}
            for row in dyn_rows:
                y = int(row.get("exam__year") or 0)
                bucket = by_year.setdefault(
                    y, {"student_ids": set(), "results": 0, "sum_score": 0.0, "passed": 0}
                )
                score = float(row.get("score") or 0)
                sid = row.get("student_id")
                if sid is not None:
                    bucket["student_ids"].add(sid)
                bucket["results"] += 1
                bucket["sum_score"] += score
                if _is_passed_ege(row):
                    bucket["passed"] += 1
            for y in sorted(by_year.keys()):
                results_n = int(by_year[y]["results"] or 0)
                participants_n = len(by_year[y]["student_ids"]) or results_n
                dynamics.append(
                    {
                        "year": y,
                        "avg": round((by_year[y]["sum_score"] / results_n), 2) if results_n else 0.0,
                        "pass_rate": round((int(by_year[y]["passed"] or 0) / results_n) * 100, 1) if results_n else 0.0,
                        "participants": participants_n,
                        "results": results_n,
                    }
                )
        else:
            dyn_qs = (
                ExamResult.objects.filter(student__school_id=school_id, exam__exam_type=et, exam__year__in=years)
                .values("exam__year")
                .annotate(
                    avg=Avg("score"),
                    participants=Count("student_id", distinct=True),
                    results=Count("id"),
                    passed=Count("id", filter=Q(passed=True)),
                )
                .order_by("exam__year")
            )
            for row in dyn_qs:
                results_n = int(row.get("results") or 0) or int(row["participants"] or 0)
                dynamics.append(
                    {
                        "year": int(row["exam__year"]),
                        "avg": round(float(row["avg"] or 0), 2),
                        "pass_rate": round((int(row["passed"] or 0) / results_n) * 100, 1) if results_n else 0.0,
                        "participants": int(row["participants"] or 0),
                        "results": results_n,
                    }
                )

    return {
        "has_data": True,
        "generated_at": date.today().strftime("%d.%m.%Y"),
        "exam_type": et,
        "year": year,
        "total": total,
        "participants": participants,
        "subjects_count": subjects_count,
        "avg_score": round(avg_score, 2),
        "min_score": round(min_score, 2),
        "max_score": round(max_score, 2),
        "quality_rate": quality_rate,
        "pass_rate": pass_rate,
        "high_count": high_count,
        "failed_count": failed_count,
        "prev_avg": prev_avg,
        "avg_delta": avg_delta,
        "pass_delta": pass_delta,
        "district_avg": district_avg,
        "republic_avg": republic_avg,
        "risk_count": risk_count,
        "subject_rows": subject_rows,
        "distribution": bins,
        "class_rows": class_rows,
        "task_rows": task_rows,
        "weak_subjects": weak_subjects,
        "weak_tasks": weak_tasks,
        "ai_insights": ai_insights,
        "recommendations": recommendations,
        "dynamics": dynamics,
    }


def _append_school_info_stat_report_body(doc, data: dict) -> None:
    """Разделы отчёта — в соответствии с экранной формой кабинета."""
    doc.add_heading("Общие сведения", level=1)
    summary_rows = [
        ("Количество участников", data.get("participants", data["total"])),
        ("Количество предметов", data["subjects_count"]),
        ("Средний балл", data["avg_score"]),
        ("Минимальный балл", data.get("min_score", "—")),
        ("Максимальный балл", data.get("max_score", "—")),
        ("Качество знаний, %", data["quality_rate"]),
        ("Успеваемость, %", data["pass_rate"]),
        ("Высокобалльники", data["high_count"]),
        ("Неудовлетворительные результаты", data["failed_count"]),
        ("Группа риска (обучающихся)", data.get("risk_count", "—")),
    ]
    t = doc.add_table(rows=1, cols=2)
    t.rows[0].cells[0].text, t.rows[0].cells[1].text = "Показатель", "Значение"
    for k, v in summary_rows:
        r = t.add_row().cells
        r[0].text, r[1].text = str(k), str(v)
    _style_docx_table(t, header_rows=1)

    dynamics_lines = []
    if data.get("avg_delta") is not None:
        sign = "+" if data["avg_delta"] > 0 else ""
        dynamics_lines.append(f"Изменение среднего балла к прошлому году: {sign}{data['avg_delta']}.")
    if data.get("pass_delta") is not None:
        sign = "+" if data["pass_delta"] > 0 else ""
        dynamics_lines.append(f"Изменение успеваемости к прошлому году: {sign}{data['pass_delta']} п.п.")
    if data.get("district_avg") is not None or data.get("republic_avg") is not None:
        dynamics_lines.append(
            "Сравнение среднего балла: школа {school} | район {district} | республика {republic}.".format(
                school=data["avg_score"],
                district=data.get("district_avg") if data.get("district_avg") is not None else "—",
                republic=data.get("republic_avg") if data.get("republic_avg") is not None else "—",
            )
        )
    if dynamics_lines:
        doc.add_heading("Динамика и сравнение", level=1)
        for line in dynamics_lines:
            doc.add_paragraph(line)

    if data.get("ai_insights"):
        doc.add_heading("Ключевые выводы", level=1)
        for item in data["ai_insights"]:
            doc.add_paragraph(str(item))

    doc.add_heading("Распределение по уровням", level=1)
    dist_t = doc.add_table(rows=1, cols=2)
    dist_t.rows[0].cells[0].text, dist_t.rows[0].cells[1].text = "Диапазон / оценка", "Количество"
    for label, value in data.get("distribution") or []:
        r = dist_t.add_row().cells
        r[0].text, r[1].text = str(label), str(value)
    _style_docx_table(dist_t, header_rows=1)

    if data.get("dynamics"):
        doc.add_heading("Динамика с прошлых лет", level=1)
        dyn_t = doc.add_table(rows=1, cols=4)
        dh = dyn_t.rows[0].cells
        dh[0].text, dh[1].text, dh[2].text, dh[3].text = "Год", "Результатов", "Средний", "Успеваемость, %"
        for row in data["dynamics"]:
            r = dyn_t.add_row().cells
            r[0].text = str(row["year"])
            r[1].text = str(row.get("results", row.get("participants", "")))
            r[2].text = str(row.get("avg", ""))
            r[3].text = str(row.get("pass_rate", ""))
        _style_docx_table(dyn_t, header_rows=1)

    doc.add_heading("Результаты по предметам", level=1)
    st = doc.add_table(rows=1, cols=8)
    h = st.rows[0].cells
    h[0].text, h[1].text, h[2].text, h[3].text = "Предмет", "Участ.", "Средний", "Успеваемость, %"
    h[4].text, h[5].text, h[6].text, h[7].text = "Качество, %", "Мин", "Макс", "Неуд."
    for row in data.get("subject_rows") or []:
        r = st.add_row().cells
        r[0].text = row.get("exam__subject") or "Предмет не указан"
        r[1].text = str(row.get("participants", ""))
        r[2].text = f"{float(row.get('avg') or 0):.2f}"
        r[3].text = str(row.get("pass_rate", ""))
        r[4].text = str(row.get("quality_rate", ""))
        r[5].text = f"{float(row.get('min_v') or 0):.2f}"
        r[6].text = f"{float(row.get('max_v') or 0):.2f}"
        r[7].text = str(row.get("failed", ""))
    _style_docx_table(st, header_rows=1)

    if data.get("class_rows"):
        doc.add_heading("Поклассная статистика", level=1)
        ct = doc.add_table(rows=1, cols=4)
        ch = ct.rows[0].cells
        ch[0].text, ch[1].text, ch[2].text, ch[3].text = "Класс", "Участ.", "Средний", "Успеваемость, %"
        for row in data["class_rows"]:
            r = ct.add_row().cells
            r[0].text = str(row.get("student__grade") or "Класс не указан")
            r[1].text = str(row.get("participants", ""))
            r[2].text = f"{float(row.get('avg') or 0):.2f}"
            r[3].text = str(row.get("pass_rate", ""))
        _style_docx_table(ct, header_rows=1)

    if data.get("weak_subjects"):
        doc.add_heading("Проблемные зоны", level=1)
        wt = doc.add_table(rows=1, cols=4)
        wh = wt.rows[0].cells
        wh[0].text, wh[1].text, wh[2].text, wh[3].text = "Предмет", "Средний", "Успеваемость, %", "Уровень риска"
        for row in data["weak_subjects"]:
            pr = float(row.get("pass_rate") or 0)
            risk = "Критический" if pr < 60 else "Средний" if pr < 75 else "Низкий"
            r = wt.add_row().cells
            r[0].text = row.get("exam__subject") or "Предмет не указан"
            r[1].text = f"{float(row.get('avg') or 0):.2f}"
            r[2].text = str(row.get("pass_rate", ""))
            r[3].text = risk
        _style_docx_table(wt, header_rows=1)

    if data.get("weak_tasks"):
        doc.add_heading("Анализ заданий КИМ (наиболее сложные)", level=1)
        tt = doc.add_table(rows=1, cols=5)
        th = tt.rows[0].cells
        th[0].text, th[1].text, th[2].text, th[3].text, th[4].text = (
            "Предмет",
            "Задание",
            "% выполнения",
            "Риск",
            "Анализ",
        )
        for row in data["weak_tasks"]:
            sr = float(row.get("success_rate") or 0)
            analysis = "Критическая зона" if sr < 40 else "Зона риска" if sr < 60 else "Стабильно"
            r = tt.add_row().cells
            r[0].text = row.get("exam__subject") or "Предмет"
            r[1].text = f"№{row.get('task_number', '')}"
            r[2].text = str(row.get("success_rate", ""))
            r[3].text = str(row.get("risk", ""))
            r[4].text = analysis
        _style_docx_table(tt, header_rows=1)

    if data.get("recommendations"):
        doc.add_heading("Рекомендации по результатам анализа", level=1)
        for item in data["recommendations"]:
            doc.add_paragraph(str(item), style="List Bullet")


def generate_school_info_stat_docx(school_id: int, exam_type: str, year: int | None = None) -> BytesIO:
    from docx import Document

    from users.report_ui.school_info_stat import build_info_stat_presentation
    from users.report_ui.school_info_stat_docx import render_info_stat_docx

    data = _build_school_info_stat_payload(school_id, exam_type, year)
    doc = Document()
    if not data["has_data"]:
        doc.add_paragraph(data["message"])
        output = BytesIO()
        doc.save(output)
        output.seek(0)
        return output

    ui = build_info_stat_presentation(data)
    render_info_stat_docx(doc, ui)

    output = BytesIO()
    doc.save(output)
    output.seek(0)
    return output


def generate_school_info_stat_xlsx(school_id: int, exam_type: str, year: int | None = None) -> BytesIO:
    from openpyxl import Workbook
    from openpyxl.styles import Font

    data = _build_school_info_stat_payload(school_id, exam_type, year)
    wb = Workbook()
    ws = wb.active
    ws.title = "Инфо-стат отчет"
    if not data["has_data"]:
        ws["A1"] = data["message"]
        output = BytesIO()
        wb.save(output)
        output.seek(0)
        return output

    et_label = "ЕГЭ" if data["exam_type"] == "ege" else "ОГЭ"
    year_label = f" за {data['year']} год" if data.get("year") else ""
    ws["A1"] = f"ИНФОРМАЦИОННО-СТАТИСТИЧЕСКИЙ ОТЧЕТ ГИА ({et_label}){year_label}"
    ws["A1"].font = Font(bold=True, size=14)
    ws.append(["Дата формирования", data.get("generated_at", "")])
    ws.append([])
    ws.append(["Показатель", "Значение"])
    for k, v in [
        ("Количество участников", data.get("participants", data["total"])),
        ("Количество предметов", data["subjects_count"]),
        ("Средний балл", data["avg_score"]),
        ("Минимальный балл", data.get("min_score")),
        ("Максимальный балл", data.get("max_score")),
        ("Качество знаний, %", data["quality_rate"]),
        ("Успеваемость, %", data["pass_rate"]),
        ("Высокобалльники", data["high_count"]),
        ("Неудовлетворительные результаты", data["failed_count"]),
        ("Группа риска", data.get("risk_count")),
    ]:
        ws.append([k, v])
    if data.get("avg_delta") is not None or data.get("pass_delta") is not None:
        ws.append([])
        ws.append(["Динамика и сравнение", ""])
        if data.get("avg_delta") is not None:
            ws.append(["Изменение среднего балла", data["avg_delta"]])
        if data.get("pass_delta") is not None:
            ws.append(["Изменение успеваемости, п.п.", data["pass_delta"]])
        ws.append(
            [
                "Сравнение (школа / район / республика)",
                f"{data['avg_score']} / {data.get('district_avg', '—')} / {data.get('republic_avg', '—')}",
            ]
        )
    if data.get("ai_insights"):
        ws.append([])
        ws.append(["Ключевые выводы", ""])
        for item in data["ai_insights"]:
            ws.append([str(item)])
    ws.append([])
    ws.append(["Распределение", "Количество"])
    for label, value in data.get("distribution") or []:
        ws.append([label, value])
    if data.get("dynamics"):
        ws.append([])
        ws.append(["Год", "Результатов", "Средний", "Успеваемость, %"])
        for row in data["dynamics"]:
            ws.append([row["year"], row.get("results", row.get("participants")), row.get("avg"), row.get("pass_rate")])
    ws.append([])
    ws.append(["Предмет", "Участ.", "Средний", "Успеваемость, %", "Качество, %", "Мин", "Макс", "Неуд."])
    for row in data.get("subject_rows") or []:
        ws.append([
            row.get("exam__subject") or "Предмет не указан",
            row.get("participants"),
            round(float(row.get("avg") or 0), 2),
            row.get("pass_rate"),
            row.get("quality_rate"),
            round(float(row.get("min_v") or 0), 2),
            round(float(row.get("max_v") or 0), 2),
            row.get("failed"),
        ])
    if data.get("class_rows"):
        ws.append([])
        ws.append(["Класс", "Участ.", "Средний", "Успеваемость, %"])
        for row in data["class_rows"]:
            ws.append([
                row.get("student__grade") or "Класс не указан",
                row.get("participants"),
                round(float(row.get("avg") or 0), 2),
                row.get("pass_rate"),
            ])
    if data.get("weak_subjects"):
        ws.append([])
        ws.append(["Проблемные зоны", ""])
        ws.append(["Предмет", "Средний", "Успеваемость, %"])
        for row in data["weak_subjects"]:
            ws.append([
                row.get("exam__subject"),
                round(float(row.get("avg") or 0), 2),
                row.get("pass_rate"),
            ])
    if data.get("weak_tasks"):
        ws.append([])
        ws.append(["Анализ заданий КИМ", ""])
        ws.append(["Предмет", "Задание", "% выполнения", "Риск"])
        for row in data["weak_tasks"]:
            ws.append([
                row.get("exam__subject"),
                row.get("task_number"),
                row.get("success_rate"),
                row.get("risk"),
            ])
    if data.get("recommendations"):
        ws.append([])
        ws.append(["Рекомендации", ""])
        for item in data["recommendations"]:
            ws.append([str(item)])
    output = BytesIO()
    wb.save(output)
    output.seek(0)
    return output


def generate_school_info_stat_pdf(school_id: int, exam_type: str, year: int | None = None) -> BytesIO:
    from reportlab.lib.units import mm
    from reportlab.platypus import Spacer

    data = _build_school_info_stat_payload(school_id, exam_type, year)
    font_name = _pdf_register_cyrillic_font()
    title_s, h1_s, body_s, _s = _pdf_make_styles(font_name)
    story = []
    if not data["has_data"]:
        story.append(_pdf_p(data["message"], body_s))
        return _pdf_build_document(story)

    et_label = "ЕГЭ" if data["exam_type"] == "ege" else "ОГЭ"
    year_note = f" за {data['year']} год" if data.get("year") else ""
    story.append(_pdf_p(f"ИНФОРМАЦИОННО-СТАТИСТИЧЕСКИЙ ОТЧЕТ ГИА ({et_label}){year_note}", title_s))
    story.append(_pdf_p(f"Дата формирования: {data.get('generated_at', '')}", body_s))
    story.append(Spacer(1, 6))

    story.append(_pdf_p("Сводные показатели", h1_s))
    summary = [
        ["Количество участников", str(data.get("participants", data["total"]))],
        ["Количество предметов", str(data["subjects_count"])],
        ["Средний балл", str(data["avg_score"])],
        ["Качество знаний, %", str(data["quality_rate"])],
        ["Успеваемость, %", str(data["pass_rate"])],
        ["Высокобалльники", str(data["high_count"])],
        ["Неудовлетворительные результаты", str(data["failed_count"])],
        ["Группа риска", str(data.get("risk_count", "—"))],
    ]
    story.append(_pdf_table(["Показатель", "Значение"], summary, [78 * mm, 78 * mm], font_name))
    story.append(Spacer(1, 8))

    if data.get("avg_delta") is not None or data.get("district_avg") is not None:
        story.append(_pdf_p("Динамика и сравнение", h1_s))
        if data.get("avg_delta") is not None:
            story.append(_pdf_p(f"Изменение среднего балла: {data['avg_delta']}", body_s))
        if data.get("pass_delta") is not None:
            story.append(_pdf_p(f"Изменение успеваемости: {data['pass_delta']} п.п.", body_s))
        story.append(
            _pdf_p(
                f"Сравнение: школа {data['avg_score']} | район {data.get('district_avg', '—')} | "
                f"республика {data.get('republic_avg', '—')}",
                body_s,
            )
        )
        story.append(Spacer(1, 8))

    if data.get("ai_insights"):
        story.append(_pdf_p("Ключевые выводы", h1_s))
        for item in data["ai_insights"]:
            story.append(_pdf_p(f"• {item}", body_s))
        story.append(Spacer(1, 8))

    story.append(_pdf_p("Распределение по уровням", h1_s))
    dist_rows = [[str(a), str(b)] for a, b in data.get("distribution") or []]
    story.append(_pdf_table(["Диапазон / оценка", "Количество"], dist_rows, [88 * mm, 68 * mm], font_name))
    story.append(Spacer(1, 8))

    if data.get("dynamics"):
        story.append(_pdf_p("Динамика с прошлых лет", h1_s))
        dyn_rows = [
            [str(r["year"]), str(r.get("results", r.get("participants", ""))), str(r.get("avg", "")), str(r.get("pass_rate", ""))]
            for r in data["dynamics"]
        ]
        story.append(_pdf_table(["Год", "Результатов", "Средний", "Усп., %"], dyn_rows, [28 * mm, 28 * mm, 28 * mm, 28 * mm], font_name))
        story.append(Spacer(1, 8))

    story.append(_pdf_p("Результаты по предметам", h1_s))
    subj_rows = []
    for row in data.get("subject_rows") or []:
        subj_rows.append(
            [
                str(row.get("exam__subject") or "Предмет не указан")[:28],
                str(row.get("participants", "")),
                f"{float(row.get('avg') or 0):.2f}",
                str(row.get("pass_rate", "")),
                str(row.get("quality_rate", "")),
                f"{float(row.get('min_v') or 0):.1f}",
                f"{float(row.get('max_v') or 0):.1f}",
                str(row.get("failed", "")),
            ]
        )
    story.append(
        _pdf_table(
            ["Предмет", "Уч.", "Ср.", "Усп.%", "Кач.%", "Мин", "Макс", "Неуд."],
            subj_rows,
            [34 * mm, 16 * mm, 16 * mm, 18 * mm, 18 * mm, 16 * mm, 16 * mm, 16 * mm],
            font_name,
        )
    )
    story.append(Spacer(1, 8))

    if data.get("class_rows"):
        story.append(_pdf_p("Поклассная статистика", h1_s))
        class_pdf_rows = [
            [
                str(r.get("student__grade") or "—")[:12],
                str(r.get("participants", "")),
                f"{float(r.get('avg') or 0):.2f}",
                str(r.get("pass_rate", "")),
            ]
            for r in data["class_rows"]
        ]
        story.append(_pdf_table(["Класс", "Уч.", "Средний", "Усп.%"], class_pdf_rows, [40 * mm, 28 * mm, 32 * mm, 32 * mm], font_name))
        story.append(Spacer(1, 8))

    if data.get("weak_subjects"):
        story.append(_pdf_p("Проблемные зоны", h1_s))
        weak_rows = [
            [
                str(r.get("exam__subject") or "")[:32],
                f"{float(r.get('avg') or 0):.2f}",
                str(r.get("pass_rate", "")),
            ]
            for r in data["weak_subjects"]
        ]
        story.append(_pdf_table(["Предмет", "Средний", "Усп.%"], weak_rows, [70 * mm, 30 * mm, 30 * mm], font_name))
        story.append(Spacer(1, 8))

    if data.get("weak_tasks"):
        story.append(_pdf_p("Анализ заданий КИМ", h1_s))
        task_pdf_rows = [
            [
                str(r.get("exam__subject") or "")[:24],
                f"№{r.get('task_number', '')}",
                str(r.get("success_rate", "")),
                str(r.get("risk", ""))[:16],
            ]
            for r in data["weak_tasks"]
        ]
        story.append(_pdf_table(["Предмет", "Зад.", "%", "Риск"], task_pdf_rows, [50 * mm, 22 * mm, 22 * mm, 36 * mm], font_name))
        story.append(Spacer(1, 8))

    if data.get("recommendations"):
        story.append(_pdf_p("Рекомендации", h1_s))
        for item in data["recommendations"]:
            story.append(_pdf_p(f"• {item}", body_s))

    return _pdf_build_document(story)


def _weak_subject_min_avg(exam_type: str, max_score: float) -> float:
    """Минимальный средний балл для признания предмета проблемным (при низкой, но не нулевой успеваемости)."""
    if exam_type == "oge" or max_score <= 5:
        return 3.0
    return 50.0


def _is_weak_subject_row(
    row: dict,
    *,
    exam_type: str,
    max_score: float,
    pass_rate_threshold: float = 70,
) -> bool:
    pass_rate = float(row.get("pass_rate") or 0)
    avg = float(row.get("avg") or 0)
    if pass_rate < pass_rate_threshold:
        return True
    # Предмет со 100% успеваемостью не считается проблемной зоной.
    if pass_rate >= 100:
        return False
    return avg < _weak_subject_min_avg(exam_type, max_score)


def _build_school_analytic_note_payload(school_id: int, exam_type: str, year: int | None = None) -> dict:
    from analytics.engine.attempts import filter_latest_exam_results

    et = (exam_type or "").strip().lower()
    if et not in {"ege", "oge"}:
        et = "ege"
    qs = ExamResult.objects.filter(student__school_id=school_id, exam__exam_type=et)
    if year:
        qs = qs.filter(exam__year=year)
    qs = filter_latest_exam_results(qs)
    if not qs.exists():
        return {"has_data": False, "message": "Недостаточно данных для формирования аналитической справки."}

    total = qs.count()
    participants = qs.values("student_id").distinct().count()
    subjects_count = qs.values("exam__subject").distinct().count()
    avg_score = float(qs.aggregate(v=Avg("score"))["v"] or 0)
    max_score_global = float(qs.aggregate(v=Max("score"))["v"] or 100)
    quality_threshold = 4 if max_score_global <= 5 else 60
    quality_count = qs.filter(score__gte=quality_threshold).count()
    quality_rate = round((quality_count / total) * 100, 1) if total else 0.0

    threshold_cache: dict[tuple[int, str], EgePassingThreshold | None] = {}

    def _threshold_subject_key_analytic(subject_name: str) -> str | None:
        title = (subject_name or "").strip().lower()
        if "рус" in title:
            return "russian"
        if "математика" in title and "проф" in title:
            return "math_profile"
        if "математика" in title and "баз" in title:
            return "math_basic"
        if "обществ" in title:
            return "social"
        if "информат" in title:
            return "informatics"
        if "физик" in title:
            return "physics"
        if "хими" in title:
            return "chemistry"
        if "биолог" in title:
            return "biology"
        if "истори" in title:
            return "history"
        if "литератур" in title:
            return "literature"
        if "географ" in title:
            return "geography"
        if any(lang in title for lang in ("англий", "немец", "француз", "испан", "китай", "иностран")):
            return "foreign_language"
        return None

    def _is_passed_for_analytic(row: dict) -> bool:
        from exams.passing import ege_result_passed, oge_score_passed

        if et != "ege":
            return oge_score_passed(row.get("score"), row.get("passed"))
        return ege_result_passed(
            subject_name=row.get("exam__subject"),
            year=row.get("exam__year"),
            score=row.get("score"),
            passed_flag=row.get("passed"),
            exam_code=row.get("exam__code"),
            cache=threshold_cache,
        )

    if et == "ege":
        ege_rows = list(
            qs.values(
                "student_id",
                "student__grade",
                "exam__subject",
                "exam__code",
                "exam__year",
                "score",
                "passed",
            )
        )
        pass_count = sum(1 for r in ege_rows if _is_passed_for_analytic(r))
        pass_rate = round((pass_count / total) * 100, 1) if total else 0.0

        by_subject: dict[str, dict] = {}
        for row in ege_rows:
            subject_name = row.get("exam__subject") or "Предмет не указан"
            bucket = by_subject.setdefault(
                subject_name,
                {
                    "exam__subject": subject_name,
                    "student_ids": set(),
                    "results": 0,
                    "sum_score": 0.0,
                    "min_v": None,
                    "max_v": None,
                    "passed": 0,
                },
            )
            score = float(row.get("score") or 0)
            sid = row.get("student_id")
            if sid is not None:
                bucket["student_ids"].add(sid)
            bucket["results"] += 1
            bucket["sum_score"] += score
            bucket["min_v"] = score if bucket["min_v"] is None else min(bucket["min_v"], score)
            bucket["max_v"] = score if bucket["max_v"] is None else max(bucket["max_v"], score)
            if _is_passed_for_analytic(row):
                bucket["passed"] += 1
        subject_rows = sorted(by_subject.values(), key=lambda x: x["exam__subject"])
        for row in subject_rows:
            results_n = int(row.pop("results", 0) or 0)
            student_ids = row.pop("student_ids", set()) or set()
            participants_subj = len(student_ids) or results_n
            row["participants"] = participants_subj
            row["avg"] = round((row["sum_score"] / results_n) if results_n else 0.0, 2)
            row["pass_rate"] = round((int(row["passed"] or 0) / results_n) * 100, 1) if results_n else 0.0
            row.pop("sum_score", None)
            row.pop("passed", None)

        by_class: dict[str, dict] = {}
        for row in ege_rows:
            class_name = row.get("student__grade") or "Класс не указан"
            bucket = by_class.setdefault(
                class_name,
                {"student__grade": class_name, "student_ids": set(), "results": 0, "sum_score": 0.0, "passed": 0},
            )
            score = float(row.get("score") or 0)
            sid = row.get("student_id")
            if sid is not None:
                bucket["student_ids"].add(sid)
            bucket["results"] += 1
            bucket["sum_score"] += score
            if _is_passed_for_analytic(row):
                bucket["passed"] += 1
        class_rows = list(by_class.values())
        for row in class_rows:
            results_n = int(row.pop("results", 0) or 0)
            student_ids = row.pop("student_ids", set()) or set()
            participants_cls = len(student_ids) or results_n
            row["participants"] = participants_cls
            row["avg"] = round((row["sum_score"] / results_n) if results_n else 0.0, 2)
            row["pass_rate"] = round((int(row["passed"] or 0) / results_n) * 100, 1) if results_n else 0.0
            row.pop("sum_score", None)
            row.pop("passed", None)
        class_rows.sort(key=lambda x: (-float(x.get("avg") or 0), str(x.get("student__grade") or "")))
    else:
        pass_count = qs.filter(passed=True).count()
        pass_rate = round((pass_count / total) * 100, 1) if total else 0.0
        subject_rows = list(
            qs.values("exam__subject")
            .annotate(
                participants=Count("student_id", distinct=True),
                results=Count("id"),
                avg=Avg("score"),
                passed=Count("id", filter=Q(passed=True)),
                min_v=Min("score"),
                max_v=Max("score"),
            )
            .order_by("exam__subject")
        )
        for row in subject_rows:
            results_n = int(row.pop("results", 0) or 0) or int(row.get("participants") or 0)
            row["pass_rate"] = round((int(row["passed"] or 0) / results_n) * 100, 1) if results_n else 0.0

        class_rows = list(
            qs.values("student__grade")
            .annotate(
                participants=Count("student_id", distinct=True),
                results=Count("id"),
                avg=Avg("score"),
                passed=Count("id", filter=Q(passed=True)),
            )
            .order_by("-avg", "student__grade")
        )
        for row in class_rows:
            results_n = int(row.pop("results", 0) or 0) or int(row.get("participants") or 0)
            row["pass_rate"] = round((int(row["passed"] or 0) / results_n) * 100, 1) if results_n else 0.0

    weak_subjects = sorted(
        [r for r in subject_rows if _is_weak_subject_row(r, exam_type=et, max_score=max_score_global)],
        key=lambda x: (x["pass_rate"], float(x.get("avg") or 0)),
    )

    high_threshold = 5 if max_score_global <= 5 else 70
    high_count = qs.filter(score__gte=high_threshold).count()

    if et == "ege":
        risk_students = {
            row.get("student_id") for row in ege_rows if row.get("student_id") is not None and not _is_passed_for_analytic(row)
        }
        risk_count = len(risk_students)
    else:
        risk_count = qs.filter(Q(passed=False) | Q(score__lt=3)).values("student_id").distinct().count()

    prev_avg = None
    prev_year = (year - 1) if year else None
    if prev_year:
        prev_qs = filter_latest_exam_results(
            ExamResult.objects.filter(student__school_id=school_id, exam__exam_type=et, exam__year=prev_year)
        )
        prev = prev_qs.aggregate(v=Avg("score"))["v"]
        if prev is not None:
            prev_avg = round(float(prev), 2)
    delta = round(avg_score - prev_avg, 2) if prev_avg is not None else None

    dynamics = []
    if year:
        years = [year - 2, year - 1, year]
        for y in years:
            y_qs = filter_latest_exam_results(
                ExamResult.objects.filter(student__school_id=school_id, exam__exam_type=et, exam__year=y)
            )
            if not y_qs.exists():
                continue
            results_y = y_qs.count()
            participants_y = y_qs.values("student_id").distinct().count()
            avg_y = float(y_qs.aggregate(v=Avg("score"))["v"] or 0)
            if et == "ege":
                yr_rows = list(y_qs.values("exam__subject", "exam__code", "exam__year", "score", "passed"))
                passed_y = sum(1 for r in yr_rows if _is_passed_for_analytic(r))
            else:
                passed_y = y_qs.filter(passed=True).count()
            dynamics.append(
                {
                    "year": int(y),
                    "avg": round(avg_y, 2),
                    "pass_rate": round((passed_y / results_y) * 100, 1) if results_y else 0.0,
                    "participants": participants_y,
                    "results": results_y,
                }
            )

    # Выводы и рекомендации формируются в презентационном слое
    # (факт / вывод / гипотеза / решение) без AI-домыслов.
    conclusions: list[str] = []
    recommendations: list[str] = []

    from users.report_ui.school_analytic_note_deficits import build_methodological_deficits

    methodological = build_methodological_deficits(school_id=school_id, exam_type=et, year=year)

    return {
        "has_data": True,
        "exam_type": et,
        "year": year,
        "total": total,
        "participants": participants,
        "subjects_count": subjects_count,
        "avg_score": round(avg_score, 2),
        "quality_rate": quality_rate,
        "pass_rate": pass_rate,
        "high_count": high_count,
        "risk_count": risk_count,
        "subject_rows": subject_rows,
        "class_rows": class_rows,
        "dynamics": dynamics,
        "weak_subjects": weak_subjects,
        "conclusions": conclusions,
        "recommendations": recommendations,
        "methodological": methodological,
    }


def generate_school_analytic_note_docx(school_id: int, exam_type: str, year: int | None = None) -> BytesIO:
    from docx import Document

    from users.report_ui.school_analytic_note import build_analytic_note_presentation
    from users.report_ui.school_analytic_note_docx import render_analytic_note_docx

    data = _build_school_analytic_note_payload(school_id, exam_type, year)
    doc = Document()
    if not data["has_data"]:
        doc.add_paragraph(data["message"])
        output = BytesIO()
        doc.save(output)
        output.seek(0)
        return output

    ui = build_analytic_note_presentation(data)
    render_analytic_note_docx(doc, ui)

    output = BytesIO()
    doc.save(output)
    output.seek(0)
    return output


def generate_school_analytic_note_pdf(school_id: int, exam_type: str, year: int | None = None) -> BytesIO:
    from reportlab.lib.units import mm
    from reportlab.platypus import Spacer

    data = _build_school_analytic_note_payload(school_id, exam_type, year)
    font_name = _pdf_register_cyrillic_font()
    title_s, h1_s, body_s, _s = _pdf_make_styles(font_name)
    story = []
    if not data["has_data"]:
        story.append(_pdf_p(data["message"], body_s))
        return _pdf_build_document(story)

    et_label = "ЕГЭ" if data["exam_type"] == "ege" else "ОГЭ"
    year_note = f" за {data['year']} год" if data.get("year") else ""
    story.append(_pdf_p(f"Аналитическая справка по итогам ГИА ({et_label}){year_note}", title_s))
    story.append(
        _pdf_p(
            f"Участников: {data.get('participants', data['total'])}; результатов: {data['total']}; предметов: {data['subjects_count']}; средний балл: {data['avg_score']}; "
            f"качество: {data['quality_rate']}%; успеваемость: {data['pass_rate']}%; высокобалльники: {data['high_count']}.",
            body_s,
        )
    )
    story.append(Spacer(1, 8))

    story.append(_pdf_p("По предметам", h1_s))
    subj = []
    for row in data["subject_rows"]:
        subj.append(
            [
                str(row["exam__subject"] or "—")[:32],
                str(row["participants"]),
                f"{float(row['avg'] or 0):.2f}",
                str(row["pass_rate"]),
                f"{float(row['min_v'] or 0):.2f}",
                f"{float(row['max_v'] or 0):.2f}",
            ]
        )
    story.append(
        _pdf_table(["Предмет", "Участ.", "Средний", "Усп., %", "Мин", "Макс"], subj, [44 * mm, 22 * mm, 24 * mm, 24 * mm, 22 * mm, 22 * mm], font_name)
    )
    story.append(Spacer(1, 6))

    story.append(_pdf_p("Поклассная статистика", h1_s))
    cls_rows = []
    for row in data["class_rows"]:
        cls_rows.append(
            [
                str(row.get("student__grade") or "—"),
                str(row.get("participants", "")),
                f"{float(row.get('avg') or 0):.2f}",
                str(row.get("pass_rate", "")),
            ]
        )
    story.append(_pdf_table(["Класс", "Участ.", "Средний", "Усп., %"], cls_rows, [38 * mm, 28 * mm, 32 * mm, 32 * mm], font_name))
    story.append(Spacer(1, 6))

    if data["dynamics"]:
        story.append(_pdf_p("Динамика по годам", h1_s))
        dyn = [[str(r["year"]), str(r["participants"]), str(r["avg"]), str(r["pass_rate"])] for r in data["dynamics"]]
        story.append(_pdf_table(["Год", "Участ.", "Средний", "Усп., %"], dyn, [28 * mm, 32 * mm, 36 * mm, 36 * mm], font_name))
        story.append(Spacer(1, 6))

    story.append(_pdf_p("Выводы", h1_s))
    for c in data["conclusions"][:12]:
        story.append(_pdf_p(f"• {c}", body_s))
    story.append(Spacer(1, 4))
    story.append(_pdf_p("Рекомендации", h1_s))
    for r in data["recommendations"][:12]:
        story.append(_pdf_p(f"• {r}", body_s))

    return _pdf_build_document(story)


def _topic_for_task(subject_name: str, task_number: int, exam_type: str) -> str:
    return topic_for_task(subject_name, task_number, exam_type)


def _build_school_subject_note_payload(
    school_id: int,
    exam_type: str,
    subject: str,
    year: int | None = None,
) -> dict:
    et = (exam_type or "").strip().lower()
    if et not in {"ege", "oge"}:
        et = "ege"
    subject_name = (subject or "").strip()
    if not subject_name:
        return {"has_data": False, "message": "Недостаточно данных для формирования предметной аналитической справки."}

    if not subject_name or not year:
        return {"has_data": False, "message": "Недостаточно данных для формирования предметной аналитической справки."}

    from analytics.engine.attempts import filter_latest_exam_results, task_results_for_exam_results

    base_qs = ExamResult.objects.filter(
        student__school_id=school_id,
        exam__exam_type=et,
        exam__subject=subject_name,
        exam__year=year,
    )
    qs = filter_latest_exam_results(base_qs)
    if not qs.exists():
        return {"has_data": False, "message": "Недостаточно данных для формирования предметной аналитической справки."}

    total = qs.count()
    avg_score = float(qs.aggregate(v=Avg("score"))["v"] or 0)
    min_score = float(qs.aggregate(v=Min("score"))["v"] or 0)
    max_score = float(qs.aggregate(v=Max("score"))["v"] or 0)
    pass_count = qs.filter(passed=True).count()
    pass_rate = round((pass_count / total) * 100, 1) if total else 0.0
    quality_threshold = 4 if max_score <= 5 else 60
    quality_count = qs.filter(score__gte=quality_threshold).count()
    quality_rate = round((quality_count / total) * 100, 1) if total else 0.0

    task_qs = task_results_for_exam_results(
        TaskResult.objects.filter(
            student__school_id=school_id,
            exam__exam_type=et,
            exam__subject=subject_name,
            exam__year=year,
        ),
        base_qs,
    )
    task_rows = list(
        task_qs.values("task_number")
        .annotate(
            total=Count("id"),
            plus=Count("id", filter=~Q(value__in=["-", "0", ""])),
        )
        .order_by("task_number")
    )
    for row in task_rows:
        total_answers = int(row["total"] or 0)
        plus = int(row["plus"] or 0)
        row["success_rate"] = round((plus / total_answers) * 100, 1) if total_answers else 0.0
        row["minus"] = max(total_answers - plus, 0)
        row["difficulty"] = "повышенный" if row["success_rate"] < 50 else ("базовый" if row["success_rate"] < 80 else "рабочий")
        row["analysis"] = (
            "Низкая успешность, требуется адресная отработка."
            if row["success_rate"] < 50
            else "Стабильное выполнение."
        )

    weak_tasks = [row for row in task_rows if row["success_rate"] < 50]
    strong_tasks = [row for row in task_rows if row["success_rate"] >= 80]
    critical_tasks = [row for row in task_rows if row["success_rate"] < 30]
    significant_tasks = [row for row in task_rows if 30 <= row["success_rate"] < 50]
    moderate_tasks = [row for row in task_rows if 50 <= row["success_rate"] < 70]
    weak_topics = []
    for row in weak_tasks[:8]:
        topic = _topic_for_task(subject_name, int(row["task_number"]), et)
        weak_topics.append({"task": row["task_number"], "topic": topic, "success_rate": row["success_rate"]})

    # Simple comparative dynamics by year for selected subject.
    dynamics = []
    if year:
        years = [year - 2, year - 1, year]
        dyn_qs = filter_latest_exam_results(
            ExamResult.objects.filter(
                student__school_id=school_id,
                exam__exam_type=et,
                exam__subject=subject_name,
                exam__year__in=years,
            )
        )
        for y in years:
            y_qs = dyn_qs.filter(exam__year=y)
            participants = y_qs.count()
            if not participants:
                continue
            passed = y_qs.filter(passed=True).count()
            avg_val = float(y_qs.aggregate(v=Avg("score"))["v"] or 0)
            dynamics.append(
                {
                    "year": y,
                    "avg": round(avg_val, 2),
                    "pass_rate": round((passed / participants) * 100, 1) if participants else 0.0,
                }
            )

    high_threshold = 5 if max_score <= 5 else 70
    high_count = qs.filter(score__gte=high_threshold).count()
    risk_count = qs.filter(score__lt=(3 if max_score <= 5 else 50)).count()

    conclusions = []
    quick_findings = []
    quick_findings.append(
        f"Средний результат по предмету: {round(avg_score,2)}; успеваемость: {pass_rate}%."
    )
    if weak_tasks:
        conclusions.append(
            f"Наибольшие затруднения по заданиям: {', '.join('№'+str(item['task_number']) for item in weak_tasks[:5])}."
        )
        quick_findings.append(
            f"Выявлены проблемные задания: {', '.join('№'+str(item['task_number']) for item in weak_tasks[:5])}."
        )
    else:
        quick_findings.append("Критических дефицитов по заданиям (<50%) не выявлено.")
    conclusions.append(
        f"По предмету «{subject_name}» средний балл {round(avg_score,2)}, успеваемость {pass_rate}%."
    )
    if dynamics:
        first, last = dynamics[0], dynamics[-1]
        delta = round(last["avg"] - first["avg"], 2)
        conclusions.append(
            f"Динамика {first['year']}–{last['year']}: {'рост' if delta > 0 else 'снижение' if delta < 0 else 'стабильность'} ({delta:+})."
        )
        quick_findings.append(
            f"Динамика: {'рост' if delta > 0 else 'снижение' if delta < 0 else 'стабильность'} среднего результата ({delta:+})."
        )

    deficit_summary = [
        f"Критический дефицит (<30%): {len(critical_tasks)} заданий.",
        f"Значимый дефицит (30–50%): {len(significant_tasks)} заданий.",
        f"Умеренный дефицит (50–70%): {len(moderate_tasks)} заданий.",
    ]
    if et == "ege":
        deficit_summary.append(
            "Фокус ЕГЭ: усилить задания с развернутым решением и аргументацией по критериям оценивания."
        )
    else:
        deficit_summary.append(
            "Фокус ОГЭ: стабилизировать базовые задания и поэтапно наращивать долю верных решений в заданиях повышенной сложности."
        )

    key_problems = []
    if weak_topics:
        key_problems.append("Проблемные тематические блоки: " + "; ".join(t["topic"] for t in weak_topics[:5]) + ".")
    if critical_tasks:
        key_problems.append(
            "Системные риски по заданиям: " + ", ".join(f"№{t['task_number']}" for t in critical_tasks[:8]) + "."
        )
    if not key_problems:
        key_problems.append("Системные проблемные зоны по заданным порогам не выявлены.")

    recommendations = []
    if critical_tasks or weak_tasks:
        focus = (critical_tasks or weak_tasks)[:5]
        recommendations.append(
            "Адресные консультации по заданиям: "
            + ", ".join(f"№{t['task_number']} ({t['success_rate']}%)" for t in focus)
            + "."
        )
    if weak_topics:
        recommendations.append(
            "Приоритетные темы для коррекции: " + "; ".join(t["topic"] for t in weak_topics[:3]) + "."
        )
    if risk_count:
        recommendations.append(f"Сопровождение группы риска по предмету: {risk_count} обучающихся.")
    if dynamics and len(dynamics) >= 2:
        first, last = dynamics[0], dynamics[-1]
        delta = round(float(last.get("avg") or 0) - float(first.get("avg") or 0), 2)
        recommendations.append(
            f"Учитывать динамику {first.get('year')}–{last.get('year')}: {delta:+} к среднему баллу."
        )
    if not recommendations:
        recommendations.append(
            f"Поддерживать текущий уровень по предмету «{subject_name}»: средний {round(avg_score, 2)}, успеваемость {pass_rate}%."
        )

    return {
        "has_data": True,
        "exam_type": et,
        "year": year,
        "subject": subject_name,
        "total": total,
        "avg_score": round(avg_score, 2),
        "min_score": round(min_score, 2),
        "max_score": round(max_score, 2),
        "quality_rate": quality_rate,
        "pass_rate": pass_rate,
        "high_count": high_count,
        "risk_count": risk_count,
        "task_rows": task_rows,
        "strong_tasks": strong_tasks[:10],
        "weak_tasks": weak_tasks[:15],
        "critical_tasks": critical_tasks[:15],
        "weak_topics": weak_topics,
        "quick_findings": quick_findings,
        "deficit_summary": deficit_summary,
        "key_problems": key_problems,
        "dynamics": dynamics,
        "conclusions": conclusions,
        "recommendations": recommendations,
    }


def generate_school_subject_note_docx(
    school_id: int,
    exam_type: str,
    subject: str,
    year: int | None = None,
) -> BytesIO:
    from docx import Document

    if not year:
        doc = Document()
        doc.add_paragraph("Недостаточно данных для формирования предметной аналитической справки.")
        output = BytesIO()
        doc.save(output)
        output.seek(0)
        return output

    exam_data = collect_subject_data_for_export(school_id, exam_type, subject, year)
    if exam_data is None:
        doc = Document()
        doc.add_paragraph("Недостаточно данных для формирования предметной аналитической справки.")
        output = BytesIO()
        doc.save(output)
        output.seek(0)
        return output

    return generate_word_doc(exam_data)


def generate_school_subject_note_pdf(
    school_id: int,
    exam_type: str,
    subject: str,
    year: int | None = None,
) -> BytesIO:
    from reportlab.lib.units import mm
    from reportlab.platypus import Spacer

    data = _build_school_subject_note_payload(school_id, exam_type, subject, year)
    font_name = _pdf_register_cyrillic_font()
    title_s, h1_s, body_s, small_s = _pdf_make_styles(font_name)
    story = []
    if not data["has_data"]:
        story.append(_pdf_p(data["message"], body_s))
        return _pdf_build_document(story)

    et_label = "ЕГЭ" if data["exam_type"] == "ege" else "ОГЭ"
    year_note = f" за {data['year']} год" if data.get("year") else ""
    story.append(_pdf_p(f"Предметная аналитическая справка: {data['subject']} ({et_label}){year_note}", title_s))
    story.append(
        _pdf_p(
            f"Участников: {data['total']}; средний: {data['avg_score']}; качество: {data['quality_rate']}%; "
            f"успеваемость: {data['pass_rate']}%; высокобалльники: {data['high_count']}; группа риска: {data['risk_count']}.",
            body_s,
        )
    )
    story.append(Spacer(1, 8))

    story.append(_pdf_p("Анализ заданий", h1_s))
    task_rows = []
    for row in data["task_rows"]:
        an = str(row.get("analysis", ""))[:120]
        task_rows.append([str(row["task_number"]), str(row["success_rate"]), str(row.get("difficulty", "")), an])
    story.append(
        _pdf_table(
            ["№", "% выполнения", "Уровень", "Анализ (фрагмент)"],
            task_rows,
            [14 * mm, 26 * mm, 32 * mm, 98 * mm],
            font_name,
        )
    )
    story.append(Spacer(1, 6))

    if data["dynamics"]:
        story.append(_pdf_p("Динамика по годам", h1_s))
        dyn = [[str(r["year"]), str(r["avg"]), str(r["pass_rate"])] for r in data["dynamics"]]
        story.append(_pdf_table(["Год", "Средний", "Усп., %"], dyn, [36 * mm, 40 * mm, 40 * mm], font_name))
        story.append(Spacer(1, 6))

    story.append(_pdf_p("Выводы", h1_s))
    for c in data["conclusions"][:10]:
        story.append(_pdf_p(f"• {c}", small_s))
    story.append(Spacer(1, 4))
    story.append(_pdf_p("Рекомендации", h1_s))
    for r in data["recommendations"][:10]:
        story.append(_pdf_p(f"• {r}", small_s))

    return _pdf_build_document(story)


MO_SUBJECT_GROUPS = {
    "math-mo": {
        "title": "МО учителей математики",
        "keywords": ("математ",),
    },
    "russian-mo": {
        "title": "МО учителей русского языка",
        "keywords": ("русск",),
    },
    "science-mo": {
        "title": "МО естественно-научного цикла",
        "keywords": ("физик", "хими", "биолог", "географ", "информат"),
    },
    "humanities-mo": {
        "title": "МО гуманитарного цикла",
        "keywords": ("истори", "общество", "литератур"),
    },
    "foreign-mo": {
        "title": "МО иностранных языков",
        "keywords": ("англий", "немец", "француз", "испан", "китай", "иностран"),
    },
}


def _build_mo_subject_filter(mo_key: str):
    """
    Фильтр предметов профиля МО.

    Важно: SQLite LIKE/icontains не делает case-folding для кириллицы,
    поэтому добавляем регистровые варианты ключевых слов.
    """
    meta = MO_SUBJECT_GROUPS.get(mo_key)
    if not meta:
        return None, None
    q = Q()
    for kw in meta["keywords"]:
        raw = str(kw or "").strip()
        if not raw:
            continue
        variants = {
            raw,
            raw.lower(),
            raw.upper(),
            raw.capitalize(),
            raw[:1].upper() + raw[1:].lower() if raw else raw,
        }
        for variant in variants:
            q |= Q(exam__subject__icontains=variant)
    return q, meta["title"]


def _match_mo_subjects(mo_key: str, subjects) -> tuple[list[str] | None, str | None]:
    """Точное сопоставление предметов профиля МО через Python casefold (кириллица)."""
    meta = MO_SUBJECT_GROUPS.get(mo_key)
    if not meta:
        return None, None
    keywords = tuple(str(k).casefold() for k in meta["keywords"] if k)
    matched: list[str] = []
    seen: set[str] = set()
    for subject in subjects:
        name = str(subject or "").strip()
        if not name or name in seen:
            continue
        low = name.casefold()
        if any(kw in low for kw in keywords):
            matched.append(name)
            seen.add(name)
    return matched, meta["title"]


def _build_school_mo_payload(school_id: int, exam_type: str, mo_key: str, year: int | None = None) -> dict:
    et = (exam_type or "").strip().lower()
    if et not in {"ege", "oge"}:
        et = "ege"
    if mo_key not in MO_SUBJECT_GROUPS:
        return {"has_data": False, "message": "Недостаточно данных для формирования отчета методического объединения."}
    mo_title = MO_SUBJECT_GROUPS[mo_key]["title"]

    qs = ExamResult.objects.filter(student__school_id=school_id, exam__exam_type=et)
    if year:
        qs = qs.filter(exam__year=year)
    matched, _ = _match_mo_subjects(mo_key, qs.values_list("exam__subject", flat=True).distinct())
    if not matched:
        return {"has_data": False, "message": "Недостаточно данных для формирования отчета методического объединения."}
    qs = qs.filter(exam__subject__in=matched)

    total = qs.count()
    subjects_count = qs.values("exam__subject").distinct().count()
    avg_score = float(qs.aggregate(v=Avg("score"))["v"] or 0)
    pass_count = qs.filter(passed=True).count()
    pass_rate = round((pass_count / total) * 100, 1) if total else 0.0
    max_score = float(qs.aggregate(v=Max("score"))["v"] or 100)
    quality_threshold = 4 if max_score <= 5 else 60
    quality_count = qs.filter(score__gte=quality_threshold).count()
    quality_rate = round((quality_count / total) * 100, 1) if total else 0.0
    high_threshold = 5 if max_score <= 5 else 70
    high_count = qs.filter(score__gte=high_threshold).count()
    failed_count = max(total - pass_count, 0)

    subject_rows = list(
        qs.values("exam__subject")
        .annotate(participants=Count("id"), avg=Avg("score"), passed=Count("id", filter=Q(passed=True)))
        .order_by("exam__subject")
    )
    for row in subject_rows:
        participants = int(row["participants"] or 0)
        row["pass_rate"] = round((int(row["passed"] or 0) / participants) * 100, 1) if participants else 0.0

    class_rows = list(
        qs.values("student__grade")
        .annotate(participants=Count("id"), avg=Avg("score"), passed=Count("id", filter=Q(passed=True)))
        .order_by("-avg", "student__grade")
    )
    for row in class_rows:
        participants = int(row["participants"] or 0)
        row["pass_rate"] = round((int(row["passed"] or 0) / participants) * 100, 1) if participants else 0.0

    weak_subjects = sorted(
        [r for r in subject_rows if _is_weak_subject_row(r, exam_type=et, max_score=max_score)],
        key=lambda x: (x["pass_rate"], float(x["avg"] or 0)),
    )

    dynamics = []
    if year:
        years = [year - 2, year - 1, year]
        dyn_qs = (
            ExamResult.objects.filter(
                student__school_id=school_id,
                exam__exam_type=et,
                exam__year__in=years,
                exam__subject__in=matched,
            )
            .values("exam__year")
            .annotate(avg=Avg("score"), participants=Count("id"), passed=Count("id", filter=Q(passed=True)))
            .order_by("exam__year")
        )
        for row in dyn_qs:
            participants = int(row["participants"] or 0)
            dynamics.append(
                {
                    "year": int(row["exam__year"]),
                    "avg": round(float(row["avg"] or 0), 2),
                    "pass_rate": round((int(row["passed"] or 0) / participants) * 100, 1) if participants else 0.0,
                }
            )

    conclusions = [
        f"МО: {mo_title}. Средний балл {round(avg_score, 2)}, успеваемость {pass_rate}%.",
        f"Высокобалльные результаты: {high_count}; неудовлетворительные результаты: {failed_count}.",
    ]
    if weak_subjects:
        conclusions.append(
            "Проблемные предметные зоны: "
            + ", ".join(
                f"{(r['exam__subject'] or 'предмет')} (усп. {r['pass_rate']}%)"
                for r in weak_subjects[:3]
            )
            + "."
        )
    if dynamics:
        first, last = dynamics[0], dynamics[-1]
        delta = round(last["avg"] - first["avg"], 2)
        conclusions.append(f"Динамика {first['year']}–{last['year']}: {delta:+} п. среднего балла.")

    recommendations = []
    if weak_subjects:
        recommendations.append(
            "Приоритет МО: адресная работа по предметам "
            + ", ".join((r["exam__subject"] or "предмет") for r in weak_subjects[:3])
            + f" (успеваемость ниже порога; текущий средний по МО {round(avg_score, 2)})."
        )
    if failed_count:
        recommendations.append(
            f"Разобрать {failed_count} неудовлетворительных результатов и закрепить повторную диагностику."
        )
    if dynamics and len(dynamics) >= 2:
        first, last = dynamics[0], dynamics[-1]
        delta = round(last["avg"] - first["avg"], 2)
        if delta < 0:
            recommendations.append(
                f"Остановить отрицательную динамику среднего балла ({first['year']}–{last['year']}: {delta:+})."
            )
        else:
            recommendations.append(
                f"Закрепить положительную/стабильную динамику ({first['year']}–{last['year']}: {delta:+})."
            )
    if not recommendations:
        recommendations.append(
            f"Поддерживать текущий уровень: средний {round(avg_score, 2)}, успеваемость {pass_rate}%."
        )

    work_plan = []
    for row in weak_subjects[:3]:
        subj = row.get("exam__subject") or "предмет"
        work_plan.append(
            {
                "activity": f"Консультации и разбор ошибок по предмету «{subj}» (усп. {row['pass_rate']}%)",
                "term": "ежемесячно",
                "owner": "Руководитель МО, учителя-предметники",
                "expected": f"рост успеваемости по «{subj}»",
            }
        )
    if failed_count:
        work_plan.append(
            {
                "activity": f"Индивидуальная коррекция по {failed_count} неудовлетворительным результатам",
                "term": "в течение четверти",
                "owner": "Учителя-предметники",
                "expected": "снижение доли неудовлетворительных",
            }
        )
    if dynamics:
        work_plan.append(
            {
                "activity": "Методический разбор динамики МО по загруженным протоколам",
                "term": "ежеквартально",
                "owner": "МО и администрация",
                "expected": "коррекция плана подготовки по фактическим показателям",
            }
        )

    return {
        "has_data": True,
        "mo_title": mo_title,
        "exam_type": et,
        "year": year,
        "total": total,
        "subjects_count": subjects_count,
        "avg_score": round(avg_score, 2),
        "quality_rate": quality_rate,
        "pass_rate": pass_rate,
        "high_count": high_count,
        "failed_count": failed_count,
        "subject_rows": subject_rows,
        "class_rows": class_rows,
        "dynamics": dynamics,
        "weak_subjects": weak_subjects,
        "conclusions": conclusions,
        "recommendations": recommendations,
        "work_plan": work_plan,
    }


def generate_school_mo_report_docx(school_id: int, exam_type: str, mo_key: str, year: int | None = None) -> BytesIO:
    from docx import Document

    data = _build_school_mo_payload(school_id, exam_type, mo_key, year)
    doc = Document()
    if not data["has_data"]:
        doc.add_paragraph(data["message"])
        output = BytesIO()
        doc.save(output)
        output.seek(0)
        return output

    et_label = "ЕГЭ" if data["exam_type"] == "ege" else "ОГЭ"
    year_label = f" за {data['year']} год" if data["year"] else ""
    doc.add_heading(f"Отчет методического объединения по итогам ГИА ({et_label}){year_label}", 0)
    doc.add_paragraph(data["mo_title"])

    doc.add_heading("1. Общая информация о работе МО", level=1)
    doc.add_paragraph(
        f"Количество выпускников в выборке: {data['total']}. "
        f"Количество предметов МО: {data['subjects_count']}."
    )
    doc.add_paragraph("Данные по персональному составу педагогов в текущей модели БД отсутствуют.")

    doc.add_heading("2. Анализ результатов ГИА", level=1)
    doc.add_paragraph(
        f"Средний балл: {data['avg_score']}; качество знаний: {data['quality_rate']}%; "
        f"успеваемость: {data['pass_rate']}%."
    )

    doc.add_heading("3. Показатели по предметам", level=1)
    st = doc.add_table(rows=1, cols=4)
    h = st.rows[0].cells
    h[0].text, h[1].text, h[2].text, h[3].text = "Предмет", "Участников", "Средний", "Успеваемость, %"
    for row in data["subject_rows"]:
        r = st.add_row().cells
        r[0].text = row["exam__subject"] or "Предмет не указан"
        r[1].text = str(row["participants"])
        r[2].text = f"{float(row['avg'] or 0):.2f}"
        r[3].text = str(row["pass_rate"])
    _style_docx_table(st, header_rows=1)

    doc.add_heading("4. Проблемные зоны", level=1)
    weak = data.get("weak_subjects") or []
    if weak:
        for row in weak[:6]:
            doc.add_paragraph(
                f"{row['exam__subject']}: средний балл {round(float(row['avg'] or 0),2)}, успеваемость {row['pass_rate']}%."
            )
    else:
        doc.add_paragraph("Выраженные проблемные зоны по предметам МО не выявлены.")

    doc.add_heading("5. Анализ высокобалльников", level=1)
    doc.add_paragraph(f"Количество высокобалльников: {data['high_count']}.")

    doc.add_heading("6. Сравнительный анализ и динамика", level=1)
    if data["dynamics"]:
        for row in data["dynamics"]:
            doc.add_paragraph(f"{row['year']}: средний {row['avg']}, успеваемость {row['pass_rate']}%.")
    else:
        doc.add_paragraph("Недостаточно данных для динамики.")

    doc.add_heading("7. Выводы", level=1)
    for line in data["conclusions"]:
        doc.add_paragraph(line)

    doc.add_heading("8. Рекомендации", level=1)
    for line in data["recommendations"]:
        doc.add_paragraph(line)

    doc.add_heading("9. План дальнейшей работы", level=1)
    pt = doc.add_table(rows=1, cols=4)
    ph = pt.rows[0].cells
    ph[0].text, ph[1].text, ph[2].text, ph[3].text = "Мероприятие", "Срок", "Ответственный", "Ожидаемый результат"
    work_plan = data.get("work_plan") or []
    if not work_plan:
        work_plan = [
            {
                "activity": "План формируется только при наличии проблемных показателей в протоколах",
                "term": "—",
                "owner": "—",
                "expected": "—",
            }
        ]
    for item in work_plan:
        r = pt.add_row().cells
        r[0].text = str(item.get("activity") or "—")
        r[1].text = str(item.get("term") or "—")
        r[2].text = str(item.get("owner") or "—")
        r[3].text = str(item.get("expected") or "—")
    _style_docx_table(pt, header_rows=1)

    output = BytesIO()
    doc.save(output)
    output.seek(0)
    return output


def generate_school_mo_report_pdf(school_id: int, exam_type: str, mo_key: str, year: int | None = None) -> BytesIO:
    from reportlab.lib.units import mm
    from reportlab.platypus import Spacer

    data = _build_school_mo_payload(school_id, exam_type, mo_key, year)
    font_name = _pdf_register_cyrillic_font()
    title_s, h1_s, body_s, _s = _pdf_make_styles(font_name)
    story = []
    if not data["has_data"]:
        story.append(_pdf_p(data["message"], body_s))
        return _pdf_build_document(story)

    et_label = "ЕГЭ" if data["exam_type"] == "ege" else "ОГЭ"
    year_note = f" за {data['year']} год" if data.get("year") else ""
    story.append(_pdf_p(f"Отчет методического объединения по итогам ГИА ({et_label}){year_note}", title_s))
    story.append(_pdf_p(str(data["mo_title"]), h1_s))
    story.append(
        _pdf_p(
            f"Участников: {data['total']}; предметов: {data['subjects_count']}; средний: {data['avg_score']}; "
            f"качество: {data['quality_rate']}%; успеваемость: {data['pass_rate']}%; "
            f"высокобалльники: {data['high_count']}; неудовлетворительные: {data['failed_count']}.",
            body_s,
        )
    )
    story.append(Spacer(1, 8))

    story.append(_pdf_p("Показатели по предметам", h1_s))
    rows = []
    for row in data["subject_rows"]:
        rows.append(
            [
                str(row["exam__subject"] or "—")[:36],
                str(row["participants"]),
                f"{float(row['avg'] or 0):.2f}",
                str(row["pass_rate"]),
            ]
        )
    story.append(_pdf_table(["Предмет", "Участ.", "Средний", "Усп., %"], rows, [52 * mm, 28 * mm, 32 * mm, 32 * mm], font_name))
    story.append(Spacer(1, 6))

    story.append(_pdf_p("Поклассно", h1_s))
    cls = []
    for row in data["class_rows"]:
        cls.append(
            [
                str(row.get("student__grade") or "—"),
                str(row.get("participants", "")),
                f"{float(row.get('avg') or 0):.2f}",
                str(row.get("pass_rate", "")),
            ]
        )
    story.append(_pdf_table(["Класс", "Участ.", "Средний", "Усп., %"], cls, [38 * mm, 28 * mm, 32 * mm, 32 * mm], font_name))
    story.append(Spacer(1, 6))

    story.append(_pdf_p("Выводы", h1_s))
    for line in data["conclusions"][:8]:
        story.append(_pdf_p(f"• {line}", body_s))
    story.append(Spacer(1, 4))
    story.append(_pdf_p("Рекомендации", h1_s))
    for line in data["recommendations"][:8]:
        story.append(_pdf_p(f"• {line}", body_s))

    return _pdf_build_document(story)


def _build_school_deputy_report_payload(school_id: int, exam_type: str = "ege", year: int | None = None) -> dict:
    et = (exam_type or "").strip().lower()
    if et not in {"ege", "oge"}:
        et = "ege"
    qs = ExamResult.objects.filter(student__school_id=school_id, exam__exam_type=et)
    if year:
        qs = qs.filter(exam__year=year)
    if not qs.exists():
        label = "ЕГЭ" if et == "ege" else "ОГЭ"
        return {"has_data": False, "message": f"Недостаточно данных для формирования отчета заместителя директора по {label}."}

    def _threshold_subject_key(subject_name: str) -> str | None:
        title = (subject_name or "").strip().lower()
        if "рус" in title:
            return "russian"
        if "СЂСѓСЃ" in title:
            return "russian"
        if "математика" in title and "проф" in title:
            return "math_profile"
        if "математика" in title and "баз" in title:
            return "math_basic"
        if "обществ" in title:
            return "social"
        if "кегэ" in title:
            return "informatics"
        if "информат" in title:
            return "informatics"
        if "физик" in title:
            return "physics"
        if "хими" in title:
            return "chemistry"
        if "биолог" in title:
            return "biology"
        if "истори" in title:
            return "history"
        if "литератур" in title:
            return "literature"
        if "географ" in title:
            return "geography"
        if any(lang in title for lang in ("англий", "немец", "француз", "испан", "китай", "иностран")):
            return "foreign_language"
        return None

    threshold_cache: dict[tuple[int, str], EgePassingThreshold | None] = {}

    def _is_passed(row: dict) -> bool:
        from exams.passing import ege_result_passed, oge_score_passed

        if et != "ege":
            return oge_score_passed(row.get("score"), row.get("passed"))
        return ege_result_passed(
            subject_name=row.get("exam__subject"),
            year=row.get("exam__year"),
            score=row.get("score"),
            passed_flag=row.get("passed"),
            exam_code=row.get("exam__code"),
            cache=threshold_cache,
        )

    rows = list(
        qs.values(
            "student_id",
            "student__grade",
            "exam__subject",
            "exam__code",
            "exam__year",
            "score",
            "passed",
        )
    )
    total = len(rows)
    avg_score = round(sum(float(r.get("score") or 0) for r in rows) / total, 2) if total else 0.0
    pass_count = sum(1 for r in rows if _is_passed(r))
    pass_rate = round((pass_count / total) * 100, 1) if total else 0.0
    quality_threshold = 60 if et == "ege" else 4
    high_threshold = 70 if et == "ege" else 5
    quality_count = sum(1 for r in rows if float(r.get("score") or 0) >= quality_threshold)
    quality_rate = round((quality_count / total) * 100, 1) if total else 0.0
    high_count = sum(1 for r in rows if float(r.get("score") or 0) >= high_threshold)
    failed_count = max(total - pass_count, 0)
    participants = len({r.get("student_id") for r in rows if r.get("student_id") is not None})

    by_subject: dict[str, dict] = {}
    by_class: dict[str, dict] = {}
    risk_students = set()
    for r in rows:
        subject = r.get("exam__subject") or "Предмет не указан"
        cls = r.get("student__grade") or "Класс не указан"
        score = float(r.get("score") or 0)
        passed = _is_passed(r)
        if not passed:
            risk_students.add(r.get("student_id"))

        s = by_subject.setdefault(subject, {"exam__subject": subject, "participants": 0, "sum_score": 0.0, "passed": 0, "min_v": None, "max_v": None})
        s["participants"] += 1
        s["sum_score"] += score
        s["min_v"] = score if s["min_v"] is None else min(s["min_v"], score)
        s["max_v"] = score if s["max_v"] is None else max(s["max_v"], score)
        if passed:
            s["passed"] += 1

        c = by_class.setdefault(cls, {"student__grade": cls, "participants": 0, "sum_score": 0.0, "passed": 0})
        c["participants"] += 1
        c["sum_score"] += score
        if passed:
            c["passed"] += 1

    subject_rows = []
    for row in by_subject.values():
        p = int(row["participants"] or 0)
        avg = (row["sum_score"] / p) if p else 0.0
        pr = round((int(row["passed"] or 0) / p) * 100, 1) if p else 0.0
        subject_rows.append(
            {
                **row,
                "avg": round(avg, 2),
                "pass_rate": pr,
                "failed": max(p - int(row["passed"] or 0), 0),
                "risk": _risk_level_ru_from_rates(float(pr), low=60, mid=75),
            }
        )
    subject_rows = sorted(subject_rows, key=lambda x: (x["pass_rate"], -float(x["avg"] or 0), x["exam__subject"]))

    class_rows = []
    for row in by_class.values():
        p = int(row["participants"] or 0)
        avg = (row["sum_score"] / p) if p else 0.0
        pr = round((int(row["passed"] or 0) / p) * 100, 1) if p else 0.0
        class_rows.append({**row, "avg": round(avg, 2), "pass_rate": pr})
    class_rows = sorted(class_rows, key=lambda x: (-float(x["avg"] or 0), x["student__grade"]))

    selected_year = year if year else max(int(r.get("exam__year") or 0) for r in rows)
    prev_qs = ExamResult.objects.filter(student__school_id=school_id, exam__exam_type=et, exam__year=selected_year - 1)
    prev_avg = round(float(prev_qs.aggregate(v=Avg("score"))["v"] or 0), 2) if prev_qs.exists() else None
    avg_delta = round(avg_score - prev_avg, 2) if prev_avg is not None else None

    school = School.objects.filter(id=school_id).select_related("district").first()
    district_avg = None
    republic_avg = None
    if school and school.district_id:
        district_qs = ExamResult.objects.filter(student__school__district_id=school.district_id, exam__exam_type=et, exam__year=selected_year)
        if district_qs.exists():
            district_avg = round(float(district_qs.aggregate(v=Avg("score"))["v"] or 0), 2)
    republic_qs = ExamResult.objects.filter(exam__exam_type=et, exam__year=selected_year)
    if republic_qs.exists():
        republic_avg = round(float(republic_qs.aggregate(v=Avg("score"))["v"] or 0), 2)

    dynamics = []
    for dyn_year in [selected_year - 2, selected_year - 1, selected_year]:
        y_qs = ExamResult.objects.filter(
            student__school_id=school_id,
            exam__exam_type=et,
            exam__year=dyn_year,
        )
        if not y_qs.exists():
            continue
        participants_d = y_qs.count()
        avg_d = round(float(y_qs.aggregate(v=Avg("score"))["v"] or 0), 2)
        if et == "ege":
            yr_dyn = list(y_qs.values("exam__subject", "exam__code", "exam__year", "score", "passed"))
            passed_d = sum(1 for r in yr_dyn if _is_passed(r))
        else:
            passed_d = y_qs.filter(passed=True).count()
        pass_rate_d = round((passed_d / participants_d) * 100, 1) if participants_d else 0.0
        dynamics.append(
            {
                "year": int(dyn_year),
                "avg": avg_d,
                "participants": participants_d,
                "pass_rate": pass_rate_d,
            }
        )

    task_qs = TaskResult.objects.filter(student__school_id=school_id, exam__exam_type=et, exam__year=selected_year)
    task_rows = list(
        task_qs.values("exam__subject", "task_number")
        .annotate(total=Count("id"), plus=Count("id", filter=~Q(value__in=["-", "0", ""])))
        .order_by("exam__subject", "task_number")
    )
    for row in task_rows:
        total_answers = int(row["total"] or 0)
        plus = int(row["plus"] or 0)
        row["success_rate"] = round((plus / total_answers) * 100, 1) if total_answers else 0.0
        row["risk"] = _risk_level_ru_from_rates(float(row["success_rate"] or 0))
    weak_tasks = sorted(task_rows, key=lambda x: x["success_rate"])[:15]

    mo_rows = []
    for key, meta in MO_SUBJECT_GROUPS.items():
        mo_q, _ = _build_mo_subject_filter(key)
        mo_qs = qs.filter(mo_q) if mo_q else ExamResult.objects.none()
        if not mo_qs.exists():
            continue
        mo_total = mo_qs.count()
        mo_avg = round(float(mo_qs.aggregate(v=Avg("score"))["v"] or 0), 2)
        mo_pass = 0
        for row in mo_qs.values("exam__subject", "exam__year", "score", "passed"):
            if _is_passed(row):
                mo_pass += 1
        mo_rows.append(
            {
                "mo_key": key,
                "mo_title": meta["title"],
                "participants": mo_total,
                "avg": mo_avg,
                "pass_rate": round((mo_pass / mo_total) * 100, 1) if mo_total else 0.0,
            }
        )
    mo_rows = sorted(mo_rows, key=lambda x: (x["pass_rate"], -x["avg"]))

    executive_summary = []
    if avg_delta is not None:
        executive_summary.append(
            f"Динамика среднего балла к {selected_year - 1} году: {'+' if avg_delta > 0 else ''}{avg_delta}."
        )
    if subject_rows:
        executive_summary.append("Проблемные предметы: " + ", ".join(r["exam__subject"] for r in subject_rows[:3]) + ".")
    if weak_tasks:
        executive_summary.append(
            "Критические задания КИМ: " + ", ".join(f"{r['exam__subject']} №{r['task_number']}" for r in weak_tasks[:5]) + "."
        )
    executive_summary.append(f"Группа риска: {len([s for s in risk_students if s is not None])} обучающихся.")

    recommendations = []
    if subject_rows:
        recommendations.append(
            "Администрации: мониторинг предметов риска — "
            + ", ".join(
                f"{r['exam__subject']} (усп. {r['pass_rate']}%)"
                for r in subject_rows[:4]
            )
            + "."
        )
    if weak_tasks:
        recommendations.append(
            "МО: практикумы по заданиям КИМ с низкой успешностью — "
            + ", ".join(
                f"{r['exam__subject']} №{r['task_number']} ({r['success_rate']}%)"
                for r in weak_tasks[:5]
            )
            + "."
        )
    risk_n = len([s for s in risk_students if s is not None])
    if risk_n:
        recommendations.append(
            f"Учителям и классным руководителям: адресное сопровождение группы риска ({risk_n} обучающихся)."
        )
    if avg_delta is not None and avg_delta < 0:
        recommendations.append(
            f"Зафиксировать меры по восстановлению среднего балла (динамика к прошлому году: {avg_delta:+})."
        )
    if not recommendations:
        recommendations.append(
            f"Поддерживать текущий уровень: средний {avg_score}, успеваемость {pass_rate}%."
        )

    improvement_plan = []
    for row in subject_rows[:4]:
        subj = row.get("exam__subject") or "предмет"
        improvement_plan.append(
            {
                "activity": f"Диагностический срез и коррекция по «{subj}» (усп. {row['pass_rate']}%, ср. {row['avg']})",
                "term": "в течение четверти",
                "owner": "Заместитель директора, МО",
                "expected": f"рост успеваемости по «{subj}»",
            }
        )
    for row in weak_tasks[:3]:
        improvement_plan.append(
            {
                "activity": (
                    f"Консультации по {row.get('exam__subject')} №{row.get('task_number')} "
                    f"(успешность {row.get('success_rate')}%)"
                ),
                "term": "еженедельно",
                "owner": "Учителя-предметники",
                "expected": "рост успешности задания",
            }
        )
    if risk_n:
        improvement_plan.append(
            {
                "activity": f"Индивидуальные маршруты для группы риска ({risk_n} обучающихся)",
                "term": "постоянно до следующей волны ГИА",
                "owner": "Классные руководители, учителя",
                "expected": "снижение доли неудовлетворительных результатов",
            }
        )

    # Без GigaChat: executive_summary и рекомендации только из фактических метрик БД.
    school_name = school.name if school else f"ОО {school_id}"
    return {
        "has_data": True,
        "exam_type": et,
        "year": selected_year,
        "school_name": school_name,
        "generated_at": date.today().strftime("%d.%m.%Y"),
        "participants": participants,
        "total": total,
        "avg_score": avg_score,
        "quality_rate": quality_rate,
        "pass_rate": pass_rate,
        "high_count": high_count,
        "failed_count": failed_count,
        "avg_delta": avg_delta,
        "district_avg": district_avg,
        "republic_avg": republic_avg,
        "risk_count": risk_n,
        "subject_rows": subject_rows,
        "class_rows": class_rows,
        "mo_rows": mo_rows,
        "task_rows": task_rows,
        "weak_tasks": weak_tasks,
        "dynamics": dynamics,
        "executive_summary": executive_summary,
        "recommendations": recommendations,
        "improvement_plan": improvement_plan,
    }


def generate_school_deputy_report_docx(school_id: int, exam_type: str = "ege", year: int | None = None) -> BytesIO:
    from docx import Document

    data = _build_school_deputy_report_payload(school_id, exam_type, year)
    doc = Document()
    if not data["has_data"]:
        doc.add_paragraph(data["message"])
        output = BytesIO()
        doc.save(output)
        output.seek(0)
        return output

    exam_label = "ЕГЭ" if data["exam_type"] == "ege" else "ОГЭ"
    doc.add_heading(f"Отчет заместителя директора по итогам {exam_label}", 0)
    doc.add_paragraph(f"{data['school_name']}; период: {data['year']} год; дата: {data['generated_at']}.")
    doc.add_heading("1. Краткое управленческое резюме", level=1)
    for item in data["executive_summary"]:
        doc.add_paragraph(item)
    doc.add_heading(f"2. Общие результаты {exam_label}", level=1)
    doc.add_paragraph(
        f"Участников: {data['participants']}; средний балл: {data['avg_score']}; качество: {data['quality_rate']}%; "
        f"успеваемость: {data['pass_rate']}%; высокобалльники: {data['high_count']}; неудовлетворительные: {data['failed_count']}."
    )
    doc.add_heading("3. Анализ по предметам", level=1)
    t = doc.add_table(rows=1, cols=6)
    h = t.rows[0].cells
    h[0].text, h[1].text, h[2].text, h[3].text, h[4].text, h[5].text = "Предмет", "Участ.", "Средний", "Усп., %", "Мин", "Макс"
    for row in data["subject_rows"]:
        r = t.add_row().cells
        r[0].text = row["exam__subject"]
        r[1].text = str(row["participants"])
        r[2].text = str(row["avg"])
        r[3].text = str(row["pass_rate"])
        r[4].text = str(round(float(row["min_v"] or 0), 2))
        r[5].text = str(round(float(row["max_v"] or 0), 2))
    _style_docx_table(t, header_rows=1)
    doc.add_heading("4. Анализ работы МО", level=1)
    for row in data["mo_rows"]:
        doc.add_paragraph(f"{row['mo_title']}: средний {row['avg']}, успеваемость {row['pass_rate']}%, участников {row['participants']}.")
    doc.add_heading("5. Анализ КИМ и проблемные зоны", level=1)
    for row in data["weak_tasks"][:10]:
        doc.add_paragraph(f"{row['exam__subject']} №{row['task_number']}: {row['success_rate']}% ({row['risk']}).")
    doc.add_heading("6. Рекомендации", level=1)
    for item in data["recommendations"]:
        doc.add_paragraph(item)
    doc.add_heading("7. План улучшения результатов", level=1)
    pt = doc.add_table(rows=1, cols=4)
    ph = pt.rows[0].cells
    ph[0].text, ph[1].text, ph[2].text, ph[3].text = "Мероприятие", "Срок", "Ответственный", "Ожидаемый результат"
    for item in data["improvement_plan"]:
        r = pt.add_row().cells
        r[0].text, r[1].text, r[2].text, r[3].text = item["activity"], item["term"], item["owner"], item["expected"]
    _style_docx_table(pt, header_rows=1)
    output = BytesIO()
    doc.save(output)
    output.seek(0)
    return output


def generate_school_deputy_report_pdf(school_id: int, exam_type: str = "ege", year: int | None = None) -> BytesIO:
    from reportlab.lib.units import mm
    from reportlab.platypus import Spacer

    data = _build_school_deputy_report_payload(school_id, exam_type, year)
    font_name = _pdf_register_cyrillic_font()
    title_s, h1_s, body_s, small_s = _pdf_make_styles(font_name)
    story = []
    if not data["has_data"]:
        story.append(_pdf_p(data["message"], body_s))
        return _pdf_build_document(story)

    exam_label = "ЕГЭ" if data["exam_type"] == "ege" else "ОГЭ"
    story.append(_pdf_p(f"Отчет заместителя директора по итогам {exam_label}", title_s))
    story.append(_pdf_p(f"{data['school_name']} | {data['year']} г. | {data['generated_at']}", body_s))
    story.append(Spacer(1, 6))

    story.append(_pdf_p("Ключевые показатели", h1_s))
    kpi = [
        ["Участников", str(data["participants"])],
        ["Средний балл", str(data["avg_score"])],
        ["Качество, %", str(data["quality_rate"])],
        ["Успеваемость, %", str(data["pass_rate"])],
        ["Высокобалльники", str(data["high_count"])],
        ["Неудовлетворительные", str(data["failed_count"])],
        ["Группа риска", str(data["risk_count"])],
    ]
    if data.get("avg_delta") is not None:
        kpi.append(["Изменение к прошлому году", str(data["avg_delta"])])
    story.append(_pdf_table(["Показатель", "Значение"], kpi, [78 * mm, 78 * mm], font_name))
    story.append(Spacer(1, 6))

    story.append(_pdf_p("Краткое управленческое резюме", h1_s))
    for item in data["executive_summary"][:6]:
        story.append(_pdf_p(f"• {item}", small_s))
    story.append(Spacer(1, 6))

    story.append(_pdf_p("По предметам", h1_s))
    subj = []
    for row in data["subject_rows"]:
        subj.append(
            [
                str(row["exam__subject"] or "—")[:28],
                str(row["participants"]),
                str(row["avg"]),
                str(row["pass_rate"]),
                str(round(float(row["min_v"] or 0), 1)),
                str(round(float(row["max_v"] or 0), 1)),
            ]
        )
    story.append(
        _pdf_table(
            ["Предмет", "Участ.", "Средний", "Усп., %", "Мин", "Макс"],
            subj,
            [40 * mm, 22 * mm, 24 * mm, 24 * mm, 20 * mm, 20 * mm],
            font_name,
        )
    )
    story.append(Spacer(1, 6))

    story.append(_pdf_p("Поклассно", h1_s))
    cls = []
    for row in data["class_rows"]:
        cls.append(
            [
                str(row.get("student__grade") or "—"),
                str(row.get("participants", "")),
                str(round(float(row.get("avg") or 0), 2)),
                str(row.get("pass_rate", "")),
            ]
        )
    story.append(_pdf_table(["Класс", "Участ.", "Средний", "Усп., %"], cls, [38 * mm, 28 * mm, 32 * mm, 32 * mm], font_name))
    story.append(Spacer(1, 6))

    story.append(_pdf_p("Проблемные задания КИМ (фрагмент)", h1_s))
    wt = []
    for row in data["weak_tasks"][:20]:
        wt.append(
            [
                str(row.get("exam__subject", ""))[:20],
                str(row.get("task_number", "")),
                str(row.get("success_rate", "")),
                str(row.get("risk", ""))[:12],
            ]
        )
    if wt:
        story.append(_pdf_table(["Предмет", "№", "Успешн., %", "Риск"], wt, [44 * mm, 16 * mm, 28 * mm, 28 * mm], font_name))
    else:
        story.append(_pdf_p("Критические задания по порогу не выявлены.", body_s))
    story.append(Spacer(1, 6))

    story.append(_pdf_p("Рекомендации", h1_s))
    for item in data["recommendations"][:8]:
        story.append(_pdf_p(f"• {item}", small_s))

    return _pdf_build_document(story)


def generate_school_deputy_report_xlsx(school_id: int, exam_type: str = "ege", year: int | None = None) -> BytesIO:
    from openpyxl import Workbook

    data = _build_school_deputy_report_payload(school_id, exam_type, year)
    wb = Workbook()
    ws = wb.active
    exam_label = "ЕГЭ" if data.get("exam_type") == "ege" else "ОГЭ"
    ws.title = f"Отчет замдиректора {exam_label}"
    if not data["has_data"]:
        ws.append([data["message"]])
    else:
        ws.append(["Показатель", "Значение"])
        for item in [
            ("Школа", data["school_name"]),
            ("Год", data["year"]),
            ("Дата", data["generated_at"]),
            ("Участники", data["participants"]),
            ("Средний балл", data["avg_score"]),
            ("Качество, %", data["quality_rate"]),
            ("Успеваемость, %", data["pass_rate"]),
            ("Высокобалльники", data["high_count"]),
            ("Неудовлетворительные", data["failed_count"]),
            ("Группа риска", data["risk_count"]),
        ]:
            ws.append(list(item))

        ws2 = wb.create_sheet("Предметы")
        ws2.append(["Предмет", "Участ.", "Средний", "Усп., %", "Мин", "Макс", "Риск"])
        for row in data["subject_rows"]:
            ws2.append([row["exam__subject"], row["participants"], row["avg"], row["pass_rate"], row["min_v"], row["max_v"], row["risk"]])

        ws3 = wb.create_sheet("Классы")
        ws3.append(["Класс", "Участ.", "Средний", "Усп., %"])
        for row in data["class_rows"]:
            ws3.append([row["student__grade"], row["participants"], row["avg"], row["pass_rate"]])

        ws4 = wb.create_sheet("План")
        ws4.append(["Мероприятие", "Срок", "Ответственный", "Ожидаемый результат"])
        for row in data["improvement_plan"]:
            ws4.append([row["activity"], row["term"], row["owner"], row["expected"]])

    output = BytesIO()
    wb.save(output)
    output.seek(0)
    return output


from users.district_export_reports import (  # noqa: E402
    _build_district_gia_summary_payload,
    generate_district_gia_summary_docx,
)


def generate_district_gia_summary_pdf(district_id: int, exam_type: str, year: int | None = None) -> BytesIO:
    """PDF presentation for district GIA summary — formatting only."""
    from users.report_ui.district_gia_summary_pdf import render_district_gia_summary_pdf

    data = _build_district_gia_summary_payload(district_id, exam_type, year)
    return render_district_gia_summary_pdf(
        data,
        helpers={
            "register_font": _pdf_register_cyrillic_font,
            "table": _pdf_table,
            "build": _pdf_build_document,
        },
    )


def generate_district_gia_summary_xlsx(district_id: int, exam_type: str, year: int | None = None) -> BytesIO:
    from openpyxl import Workbook

    data = _build_district_gia_summary_payload(district_id, exam_type, year)
    wb = Workbook()
    ws = wb.active
    ws.title = "Свод по району"
    if not data["has_data"]:
        ws.append([data["message"]])
    else:
        ws.append(["Показатель", "Значение"])
        for item in [
            ("Год", data["year"]),
            ("Участники", data["participants"]),
            ("Средний результат", data["avg_score"]),
            ("Качество, %", data["quality_rate"]),
            ("Успеваемость, %", data["pass_rate"]),
            ("High performers", data["high_count"]),
            ("Неудовлетворительные", data["failed_count"]),
        ]:
            ws.append(list(item))

        ws2 = wb.create_sheet("По предметам")
        ws2.append(["Предмет", "Участ.", "Средний", "Усп., %", "Качество, %", "Min", "Max", "Fail rate, %"])
        for row in data["subject_rows"]:
            ws2.append([
                row["exam__subject"],
                row["participants"],
                row["avg"],
                row["pass_rate"],
                row["quality_rate"],
                row["min_v"],
                row["max_v"],
                row["fail_rate"],
            ])

        ws3 = wb.create_sheet("По школам")
        ws3.append(["Школа", "Код", "Участ.", "Средний", "Усп., %"])
        for row in data["school_rows"]:
            ws3.append([row["student__school__name"], row["student__school__code"], row["participants"], row["avg"], row["pass_rate"]])

    output = BytesIO()
    wb.save(output)
    output.seek(0)
    return output
